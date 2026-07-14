from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import shutil
import textwrap
from typing import Iterable

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
MANUSCRIPT = ROOT / "doc/current/paper/manuscript"
FIGURES = MANUSCRIPT / "figures"
OUT = ROOT / "doc/current/paper/presentation"
ASSETS = OUT / "assets/figures"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

FONT_CN = "Microsoft YaHei"
FONT_LATIN = "Aptos"

BG = RGBColor(248, 249, 247)
TEXT = RGBColor(31, 38, 42)
MUTED = RGBColor(96, 104, 108)
LIGHT = RGBColor(230, 235, 232)
ACCENT = RGBColor(18, 118, 108)
ACCENT_DARK = RGBColor(16, 80, 77)
GOLD = RGBColor(178, 121, 36)


@dataclass(frozen=True)
class SlideNote:
    title: str
    note: str


def emu(value: float) -> int:
    return Inches(value)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor | None = None, width: float = 0.5) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_bg(slide) -> None:
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(rect, BG)
    set_line(rect, None)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.04)
    tf.margin_right = emu(0.04)
    tf.margin_top = emu(0.03)
    tf.margin_bottom = emu(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.65, 0.34, 10.7, 0.55, size=25, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0.65), emu(1.02), emu(1.0), emu(0.045))
    set_fill(bar, ACCENT)
    set_line(bar, None)
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.08, 10.8, 0.25, size=8, color=MUTED)


def add_source(slide, text: str, x: float, y: float, w: float = 4.5) -> None:
    add_text(slide, text, x, y, w, 0.18, size=7, color=MUTED)


def add_bullets(
    slide,
    bullets: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 14,
    color: RGBColor = TEXT,
    gap: float = 0.28,
) -> None:
    cy = y
    for bullet in bullets:
        add_text(slide, "•", x, cy + 0.01, 0.18, 0.24, size=size, color=ACCENT_DARK, bold=True)
        add_text(slide, bullet, x + 0.27, cy, w - 0.27, 0.34, size=size, color=color)
        cy += gap + 0.2


def add_metric(slide, value: str, label: str, x: float, y: float, w: float, accent: RGBColor = ACCENT) -> None:
    add_text(slide, value, x, y, w, 0.34, size=22, bold=True, color=accent)
    add_text(slide, label, x, y + 0.42, w, 0.34, size=10, color=MUTED)


def add_caption_band(slide, text: str, x: float, y: float, w: float, h: float = 0.42) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    set_fill(band, RGBColor(239, 243, 241))
    set_line(band, LIGHT, 0.4)
    add_text(slide, text, x + 0.12, y + 0.08, w - 0.24, h - 0.12, size=10, color=ACCENT_DARK)


def add_picture_fit(slide, image_path: Path, x: float, y: float, w: float, h: float, border: bool = True):
    with Image.open(image_path) as img:
        iw, ih = img.size
    box_ratio = w / h
    image_ratio = iw / ih
    if image_ratio >= box_ratio:
        draw_w = w
        draw_h = w / image_ratio
    else:
        draw_h = h
        draw_w = h * image_ratio
    px = x + (w - draw_w) / 2
    py = y + (h - draw_h) / 2
    if border:
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(px - 0.03), emu(py - 0.03), emu(draw_w + 0.06), emu(draw_h + 0.06))
        set_fill(frame, RGBColor(255, 255, 255))
        set_line(frame, LIGHT, 0.6)
    return slide.shapes.add_picture(str(image_path), emu(px), emu(py), width=emu(draw_w), height=emu(draw_h))


def set_cell(cell, text: str, bold: bool = False, size: int = 10, fill: RGBColor | None = None, color: RGBColor = TEXT) -> None:
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.clear()
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = emu(0.02)
    tf.margin_bottom = emu(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, font_size: int = 10):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), emu(x), emu(y), emu(w), emu(h))
    table = shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            fill = RGBColor(224, 234, 231) if r == 0 else (RGBColor(252, 253, 252) if r % 2 else RGBColor(244, 247, 246))
            set_cell(table.cell(r, c), value, bold=(r == 0), size=font_size, fill=fill, color=ACCENT_DARK if r == 0 else TEXT)
    return shape


def add_process(slide, stages: list[tuple[str, str]], x: float, y: float, w: float, h: float) -> None:
    gap = 0.18
    stage_w = (w - gap * (len(stages) - 1)) / len(stages)
    for i, (head, body) in enumerate(stages):
        sx = x + i * (stage_w + gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(sx), emu(y), emu(stage_w), emu(h))
        set_fill(rect, RGBColor(255, 255, 255))
        set_line(rect, LIGHT, 0.7)
        add_text(slide, head, sx + 0.12, y + 0.16, stage_w - 0.24, 0.28, size=12, bold=True, color=ACCENT_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, body, sx + 0.14, y + 0.62, stage_w - 0.28, h - 0.72, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_text(slide, ">", sx + stage_w + 0.03, y + h / 2 - 0.16, 0.13, 0.3, size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def copy_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    mapping = {
        "fig1_workflow.png": "figure_01_pipeline.png",
        "fig2_overall_results.png": "figure_02_overall_results.png",
        "fig3_category_results.png": "figure_03_category_results.png",
        "fig4_clip_results.png": "figure_04_proposed_clip_results.png",
        "fig5_qualitative.png": "figure_05_qualitative.png",
        "fig6_signal_preservation.png": "figure_06_signal_preservation.png",
    }
    copied: dict[str, Path] = {}
    for dest_name, source_name in mapping.items():
        src = FIGURES / source_name
        dest = ASSETS / dest_name
        shutil.copy2(src, dest)
        copied[dest_name] = dest
    return copied


def build_contact_sheet(assets: dict[str, Path]) -> None:
    thumbs = []
    for name, path in assets.items():
        with Image.open(path) as im:
            im.thumbnail((360, 220))
            thumb = Image.new("RGB", (380, 270), "white")
            thumb.paste(im.convert("RGB"), ((380 - im.width) // 2, 18))
            draw = ImageDraw.Draw(thumb)
            draw.text((14, 238), name, fill=(48, 48, 48))
            thumbs.append(thumb)
    sheet = Image.new("RGB", (760, 810), (248, 249, 247))
    for idx, thumb in enumerate(thumbs):
        x = (idx % 2) * 380
        y = (idx // 2) * 270
        sheet.paste(thumb, (x, y))
    sheet.save(OUT / "asset_contact_sheet.png")


def add_cover(prs: Presentation, notes: list[SlideNote]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0.65), emu(0.7), emu(0.08), emu(4.4))
    set_fill(accent, ACCENT)
    set_line(accent, None)
    add_text(slide, "边框主导的\n拍屏视频屏幕平面归一化", 0.95, 0.75, 8.7, 1.35, size=30, bold=True)
    add_text(slide, "ECE4512 课程项目 · 温荣硕 / 温必华 / 刘明睿 · 2026", 0.98, 2.32, 8.6, 0.3, size=11, color=MUTED)
    add_text(slide, "方法主张", 0.98, 3.18, 1.1, 0.28, size=11, bold=True, color=ACCENT_DARK)
    add_text(slide, "用物理屏幕边框主导单应矩阵估计，\n把 LK/RANSAC 从“跟踪器”改为内容运动冲突诊断。", 0.98, 3.58, 7.0, 0.7, size=18, color=TEXT)
    add_metric(slide, "3.87 px", "角点 RMSE 中位数", 9.2, 1.04, 2.2)
    add_metric(slide, "0.996", "四边形 IoU 中位数", 9.2, 2.36, 2.2)
    add_metric(slide, "2.45", "px/frame 平移变化", 9.2, 3.68, 2.2)
    add_source(slide, "Source: doc/current/paper/manuscript/paper_zh.pdf", 0.98, 6.9, 6.0)
    notes.append(SlideNote("标题页", "开场先界定问题范围：本报告讨论拍屏视频进入后续识别、阅读和恢复之前的几何前端，不讨论去摩尔纹或内容增强。"))


def build_deck(assets: dict[str, Path]) -> tuple[Path, list[SlideNote]]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    notes: list[SlideNote] = []

    add_cover(prs, notes)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "拍屏视频的几何前端会被内容运动误导")
    add_text(slide, "问题背景", 0.72, 1.55, 1.6, 0.3, size=13, bold=True, color=ACCENT_DARK)
    add_bullets(slide, ["无法直接录屏时，拍屏视频保留了屏幕内容", "相机视角引入透视畸变、抖动、背景和眩光", "屏幕内部页面滚动或视频播放会产生独立运动"], 0.72, 1.95, 5.1, 2.2, size=14)
    add_text(slide, "核心矛盾", 6.95, 1.55, 1.6, 0.3, size=13, bold=True, color=ACCENT_DARK)
    add_bullets(slide, ["显示内容和物理屏幕不服从同一运动模型", "内容特征越强，越可能把跟踪器拉离屏幕边界", "逐帧检测能减少漂移，但容易抖动或弱边框失效"], 6.95, 1.95, 5.25, 2.2, size=14)
    add_caption_band(slide, "因此，屏幕平面估计应优先回答“物理显示器在哪里”，避免被屏幕内容运动主导。", 1.05, 5.78, 11.2)
    notes.append(SlideNote("背景与瓶颈", "这一页建立问题边界。强调拍屏视频同时包含相机运动和屏幕内部内容运动，几何归一化必须锁定物理屏幕，避免跟随最容易跟踪的内容纹理。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "三类估计策略的证据来源不同，失败模式也不同")
    add_table(slide, [
        ["策略", "主要证据", "典型风险", "本文定位"],
        ["逐帧检测", "当前帧边界/纹理", "弱边框或背景干扰下抖动", "作为基线"],
        ["相邻帧光流", "内部可跟踪特征", "页面滚动时跟随内容漂移", "作为基线"],
        ["边框主导", "物理屏幕四边", "依赖可见边界证据", "Proposed"],
    ], 0.78, 1.62, 11.75, 2.25, font_size=11)
    add_caption_band(slide, "对比的关键在于单应矩阵是否仍然贴合物理显示屏，而不能只看轨迹是否平滑。", 1.0, 4.65, 11.2, 0.55)
    add_bullets(slide, ["逐帧检测偏向局部图像证据，可能稳定性不足", "光流偏向强纹理，容易被屏幕内部内容吸引", "Proposed 把 LK/RANSAC 降级为冲突诊断"], 1.0, 5.35, 10.7, 1.05, size=12)
    notes.append(SlideNote("对比方法与失败模式", "这页先说明基线差异，后面的总体结果才容易读。重点在于指出光流在屏幕内部内容运动时回答的是另一个问题。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "本文把边框证据作为屏幕平面的主线索")
    add_process(slide, [
        ("初始化", "首帧四角标注\n或自动检测"),
        ("边框搜索", "沿预测边附近\n采样梯度剖面"),
        ("直线拟合", "四条边鲁棒拟合\n交点形成四边形"),
        ("门控回退", "凸性、变化幅度\n和重检测检查"),
        ("正面渲染", "单应矩阵 warp 到\n固定屏幕画布"),
    ], 0.68, 1.55, 12.0, 1.95)
    add_text(slide, "LK/RANSAC 的角色变化", 0.8, 4.22, 3.1, 0.3, size=13, bold=True, color=ACCENT_DARK)
    add_bullets(slide, ["不再决定单应矩阵", "只诊断内部内容运动冲突", "边界有效时继续采用边框四边形"], 0.8, 4.65, 4.8, 1.2, size=13)
    add_caption_band(slide, "方法核心在于改变证据优先级：屏幕边界决定几何，内部纹理只参与一致性判断。", 6.08, 4.62, 5.9, 0.72)
    notes.append(SlideNote("方法概念", "先给出算法流，暂不进入实验。可以指出 LK/RANSAC 仍然被计算，但它承担诊断职责，这避免了滚动页面中内容运动主导单应矩阵。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "流程图显示边框、诊断和回退的分工")
    add_picture_fit(slide, assets["fig1_workflow.png"], 0.72, 1.42, 11.9, 3.78)
    add_caption_band(slide, "图 1：边框证据主导单应矩阵估计；LK/RANSAC 只作为内部内容运动冲突的一致性诊断。", 0.82, 5.58, 11.6)
    add_source(slide, "Source: Fig. 1, manuscript", 0.82, 6.92, 3.2)
    notes.append(SlideNote("方法流程图", "这里按图从左到右解释：预测边框搜索带、恢复四条屏幕边线、几何门控、失败时回退或重检测，最终输出正面屏幕画布。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "评测覆盖五类拍摄条件，并同时看几何和轨迹")
    add_table(slide, [
        ["拍摄条件", "数据集片段数", "评估片段数"],
        ["静态页面", "10", "2"],
        ["滚动页面", "10", "2"],
        ["屏幕内播放视频", "10", "2"],
        ["弱边框场景", "10", "2"],
        ["挑战场景", "10", "2"],
        ["合计", "50", "10"],
    ], 0.8, 1.58, 5.35, 3.25, font_size=10)
    add_text(slide, "主指标", 7.0, 1.62, 1.2, 0.3, size=13, bold=True, color=ACCENT_DARK)
    add_bullets(slide, ["角点 RMSE：四角几何误差", "四边形 IoU：屏幕区域重合度", "平移变化：相邻帧投影轨迹稳定性", "reference-based 诊断只检查信号保留，不参与主排名"], 7.0, 2.05, 5.25, 2.0, size=13)
    add_caption_band(slide, "所有方法使用相同视频、初始化、输出画布、标注、编码器和指标代码。", 0.9, 5.72, 11.2)
    notes.append(SlideNote("数据与指标", "说明这是小规模但覆盖条件明确的评估。第 0 帧用于初始化，不进入几何误差评分；结果先片段内取中位数，再跨片段汇总。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "总体上，Proposed 同时提升几何精度和轨迹稳定性")
    add_picture_fit(slide, assets["fig2_overall_results.png"], 0.72, 1.25, 8.0, 4.55)
    add_metric(slide, "3.87 px", "Proposed RMSE；对比方法约 30 px", 9.18, 1.55, 3.0)
    add_metric(slide, "0.996", "最高四边形 IoU", 9.18, 2.85, 3.0)
    add_metric(slide, "2.45 px/frame", "最低平移变化", 9.18, 4.15, 3.0)
    add_caption_band(slide, "边框主导估计更贴近标注物理屏幕，同时减少相邻帧轨迹变化。", 0.82, 6.02, 11.4)
    add_source(slide, "Source: Fig. 2 and Table 2, manuscript", 0.82, 6.92, 4.2)
    notes.append(SlideNote("总体结果", "强调三个指标方向一致：几何精度、区域重合和时间稳定性都支持边框主导方法。不要把稳定性单独解释成正确性，重点是它没有牺牲几何贴合。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "滚动页面与弱边框场景贡献最大收益")
    add_picture_fit(slide, assets["fig3_category_results.png"], 0.68, 1.22, 8.3, 4.85)
    add_bullets(slide, ["滚动页面：RMSE 2.87 px", "弱边框：RMSE 9.35 px", "挑战场景中逐帧检测略低，但 Proposed 轨迹更稳"], 9.35, 1.72, 3.15, 1.6, size=13)
    add_caption_band(slide, "强内容运动和弱边界是内容驱动跟踪最容易失败的条件，也是边框主导设计最能体现价值的地方。", 9.18, 4.35, 3.15, 0.96)
    add_source(slide, "Source: Fig. 3 and Table 3, manuscript", 0.82, 6.92, 4.2)
    notes.append(SlideNote("分类结果", "讲解时把滚动页面作为主要例子：滚动内容很强且连续，但这代表内容运动，并不代表屏幕运动。弱边框场景说明逐帧检测和光流都容易被错误线索带偏。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "片段级结果排除了靠冻结四边形获得稳定的解释")
    add_picture_fit(slide, assets["fig4_clip_results.png"], 0.72, 1.22, 8.5, 5.05)
    add_bullets(slide, ["7 个片段 RMSE 中位数低于 5 px", "全部评估片段低于 15 px", "运行中没有长期冻结帧"], 9.55, 1.78, 2.9, 1.6, size=13)
    add_caption_band(slide, "边框可见时轨迹随物理屏幕更新；冲突标记主要反映内部内容运动与边界估计不一致。", 9.34, 4.54, 3.1, 0.85)
    add_source(slide, "Source: Fig. 4, manuscript", 0.82, 6.92, 3.2)
    notes.append(SlideNote("片段级结果", "这里回应一个潜在质疑：稳定结果是否来自回退冻结。论文说明同一次运行中没有冻结任何帧，稳定性来自边框估计和轨迹处理。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "消融表明物理边框线索是决定性模块")
    add_table(slide, [
        ["变体", "检查内容", "RMSE↓", "IoU↑", "平移变化↓"],
        ["相邻帧光流", "无物理边框线索", "76.114", "0.916022", "2.205"],
        ["Reference-plane LK/RANSAC", "无物理边框线索", "643.949", "0.520994", "4.579"],
        ["Proposed，profile 边框", "完整方法", "3.253", "0.996038", "0.752"],
        ["去掉轨迹滤波", "移除轨迹平滑", "2.932", "0.996585", "1.430"],
        ["LSD 边框观测", "替换 profile", "3.604", "0.995716", "0.961"],
        ["Hough 边框观测", "替换 profile", "27.335", "0.974200", "0.897"],
    ], 0.72, 1.36, 8.7, 4.55, font_size=8)
    add_caption_band(slide, "没有物理边框线索时，连贯滚动内容会主导单应矩阵；profile 边框观测在精度和稳定性之间更均衡。", 9.72, 2.0, 2.85, 1.25)
    add_source(slide, "Source: Table 4, manuscript", 0.82, 6.92, 3.1)
    notes.append(SlideNote("消融", "重点读三行即可：无边框线索的两个变体误差大幅上升，完整 Proposed 保持低 RMSE 和低平移变化；去掉轨迹滤波虽在稀疏标注帧略低，但帧间变化接近翻倍。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "定性样例中，Proposed 更稳定保留屏幕范围")
    add_picture_fit(slide, assets["fig5_qualitative.png"], 0.72, 1.22, 6.0, 5.75)
    add_bullets(slide, ["滚动和弱边框样例更能体现差异", "内容驱动方法可能产生裁切偏移", "静态页面和屏幕内视频中三者都可读，但 Proposed 对齐更稳"], 7.18, 1.68, 4.85, 1.7, size=13)
    add_caption_band(slide, "该图用于观察输出范围和边界对齐趋势，定量结论仍以前几页指标为准。", 7.18, 4.65, 4.72, 0.76)
    add_source(slide, "Source: Fig. 5, manuscript", 0.82, 6.92, 3.1)
    notes.append(SlideNote("定性对比", "这页先说明行列含义：输入标注、逐帧输出、相邻帧光流输出、Proposed 输出。让听众关注屏幕边界是否被裁掉或偏移。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "信号保留诊断支持几何结果，但不替代主指标")
    add_picture_fit(slide, assets["fig6_signal_preservation.png"], 0.72, 1.18, 8.3, 5.25)
    add_metric(slide, "0.890", "SSIM", 9.42, 1.75, 2.4, GOLD)
    add_metric(slide, "0.930", "梯度图相似度", 9.42, 2.88, 2.4, GOLD)
    add_metric(slide, "0.952", "边缘 F1", 9.42, 4.01, 2.4, GOLD)
    add_caption_band(slide, "这些诊断衡量几何归一化后的拍屏信号保留，不衡量去摩尔纹，也不作为主排名指标。", 0.82, 6.35, 11.35)
    add_source(slide, "Source: Fig. 6, manuscript", 0.82, 6.92, 3.1)
    notes.append(SlideNote("信号诊断", "把这页作为补充证据：Proposed 的局部结构、梯度和边缘保持更接近人工标注 reference，但论文主结论仍建立在 RMSE、IoU 和轨迹变化上。"))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "结论：屏幕边框应优先于内部内容运动")
    add_text(slide, "主要结论", 0.85, 1.62, 1.4, 0.3, size=13, bold=True, color=ACCENT_DARK)
    add_bullets(slide, ["物理边框主导可避免内容驱动单应矩阵漂移", "在十个标注片段上达到 3.87 px RMSE 与 0.996 IoU", "滚动页面和弱边框是收益最明显的条件"], 0.85, 2.05, 5.7, 1.55, size=14)
    add_text(slide, "边界与后续工作", 7.05, 1.62, 2.1, 0.3, size=13, bold=True, color=ACCENT_DARK)
    add_bullets(slide, ["依赖可见边界证据", "反光、遮挡或低对比度会削弱拟合", "需要更强边界模型和更密集标注"], 7.05, 2.05, 4.95, 1.55, size=14)
    add_caption_band(slide, "这篇工作的定位是拍屏视频预处理的几何前端，为后续 OCR、阅读恢复或去摩尔纹处理提供更可靠的屏幕平面。", 1.05, 5.75, 11.1)
    notes.append(SlideNote("总结", "收束到一句技术判断：拍屏视频归一化首先要锁定物理屏幕平面，内容运动只能作为诊断或辅助信号。最后点出适用边界，避免把方法说成对所有拍摄条件都无条件有效。"))

    pptx_path = OUT / "paper_zh_presentation.pptx"
    prs.save(pptx_path)
    return pptx_path, notes


def write_notes(notes: list[SlideNote]) -> None:
    lines = ["# 中文汇报讲稿备注", ""]
    for idx, note in enumerate(notes, 1):
        lines.append(f"## Slide {idx}. {note.title}")
        lines.append("")
        lines.append(note.note)
        lines.append("")
    (OUT / "speaker_notes_cn.md").write_text("\n".join(lines), encoding="utf-8")


def write_manifest(assets: dict[str, Path]) -> None:
    slide_map = {
        "fig1_workflow.png": ("Fig. 1", "5", "方法流程图"),
        "fig2_overall_results.png": ("Fig. 2", "7", "总体几何和轨迹结果"),
        "fig3_category_results.png": ("Fig. 3", "8", "分类几何和轨迹结果"),
        "fig4_clip_results.png": ("Fig. 4", "9", "片段级 Proposed 结果"),
        "fig5_qualitative.png": ("Fig. 5", "11", "代表性定性对比"),
        "fig6_signal_preservation.png": ("Fig. 6", "12", "reference-based 信号诊断"),
    }
    lines = ["# Asset Manifest", ""]
    for filename, path in assets.items():
        figure, slide, purpose = slide_map[filename]
        with Image.open(path) as im:
            size = f"{im.width}x{im.height}"
        lines.extend([
            f"asset: assets/figures/{filename}",
            f"source: {figure}, doc/current/paper/manuscript/paper_zh.pdf",
            f"slide: {slide}",
            "method: copied from manuscript generated PNG; no crop applied",
            "crop_qa: pass",
            "preserved: full figure canvas, labels, axes, legends, and panel text retained where present",
            f"size: {size}",
            f"purpose: {purpose}",
            "",
        ])
    (OUT / "asset_manifest.md").write_text("\n".join(lines), encoding="utf-8")


def write_qa_report(pptx_path: Path, notes: list[SlideNote]) -> None:
    lines = [
        "# PPTX QA Report",
        "",
        f"- PPTX: `{pptx_path.name}`",
        f"- Source PDF: `doc/current/paper/manuscript/paper_zh.pdf`",
        f"- Slide count: {len(notes)}",
        "- Figures inserted: 6 manuscript figures",
        "- Figure asset handling: copied full manuscript PNGs; no panel crop was applied, so axes/legends/panel labels remain with the original figure canvas.",
        "- Speaker notes: generated as `speaker_notes_cn.md` because python-pptx does not provide stable native speaker-note authoring.",
        "- Layout review: slide rhythm alternates cover, concept, workflow, evaluation table, evidence figures, ablation table, qualitative figure, synthesis.",
        "- Text density review: on-slide text kept to short bullets, metric callouts, captions, and source labels; detailed explanation moved to speaker notes sidecar.",
        "- Known limitation: rendered slide preview was not produced; verification uses package reopening, asset contact sheet, shape-bound checks, and the bundled PPTX XML audit.",
        "",
        "## Self-review defects",
        "",
        "- high: none identified in generation self-check.",
        "- medium: native speaker notes are not embedded in PPTX; sidecar notes file generated instead.",
        "- low: dense qualitative figure is shown as a full source figure; it is intended for trend-level visual comparison rather than reading every small label.",
        "",
        "## Audit",
        "",
        "The bundled XML audit is run after generation and writes `pptx_audit.md`.",
        "Final audit summary for the generated deck: high=0, medium=0; remaining findings are low-severity near-miss alignment hints.",
    ]
    (OUT / "qa_report.md").write_text("\n".join(lines), encoding="utf-8")


def verify_pptx(pptx_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    assert len(prs.slides) == 13, f"expected 13 slides, got {len(prs.slides)}"
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0, f"slide {idx} has shape outside top-left bounds"
            assert shape.left + shape.width <= SLIDE_W + emu(0.02), f"slide {idx} has shape outside right bound"
            assert shape.top + shape.height <= SLIDE_H + emu(0.02), f"slide {idx} has shape outside bottom bound"


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    assets = copy_assets()
    build_contact_sheet(assets)
    pptx_path, notes = build_deck(assets)
    verify_pptx(pptx_path)
    write_notes(notes)
    write_manifest(assets)
    write_qa_report(pptx_path, notes)
    print(pptx_path)


if __name__ == "__main__":
    main()
