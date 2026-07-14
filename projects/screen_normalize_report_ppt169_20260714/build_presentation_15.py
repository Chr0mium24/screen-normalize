from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_presentation import (
    BLUE,
    BLUE_DARK,
    BLUE_PALE,
    BLUE_SOFT,
    FONT_CN,
    IMG,
    INK,
    LINE,
    MUTED,
    ORANGE,
    RED,
    RED_SOFT,
    ROOT,
    TEAL,
    TEAL_SOFT,
    WHITE,
    add_arrow,
    add_footer,
    add_header,
    add_line,
    add_notes,
    add_picture_contain,
    add_picture_crop,
    add_rect,
    add_rich_line,
    add_text,
    rgb,
    set_run_font,
    style_cell,
)


OUT = ROOT / "exports" / "screen_normalize_final_15slides_en.pptx"
W = 13.333
H = 7.5


def add_dot(slide, x, y, d=0.12, fill=BLUE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(fill)
    return shape


def add_bullet_list(slide, items, x, y, w, h, size=16, color=INK,
                    bullet_color=BLUE, gap=0.46):
    for i, item in enumerate(items):
        yy = y + i * gap
        add_dot(slide, x, yy + 0.12, 0.11, bullet_color)
        add_text(slide, item, x + 0.24, yy, w - 0.24, gap + 0.06,
                 size=size, color=color, line_spacing=1.08)


def add_metric_card(slide, x, y, w, h, value, label, accent=BLUE, fill=BLUE_SOFT,
                    value_size=25, label_size=12):
    add_rect(slide, x, y, w, h, fill=fill, line=accent, line_width=1.0)
    add_text(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.46,
             size=value_size, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.16, y + 0.74, w - 0.32, h - 0.86,
             size=label_size, color=INK, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs, prs.slide_layouts[6]


def build_deck():
    prs, blank = new_prs()

    # 1 — Cover
    slide = prs.slides.add_slide(blank)
    add_text(slide, "ECE4512 FINAL PROJECT · 2026", 0.68, 0.46, 6.8, 0.30,
             size=13.0, color=BLUE, bold=True)
    add_text(slide, "Border-Guided Screen-Plane\nNormalization for Captured-Screen Videos",
             0.68, 0.94, 7.05, 1.45, size=29.5, bold=True, line_spacing=1.02)
    add_text(slide, "Physical screen boundaries—not internal content—drive the homography",
             0.70, 2.52, 6.82, 0.42, size=16.0, color=MUTED, italic=True)
    add_rich_line(
        slide,
        [
            ("Rongshuo Wen", {"bold": True}), ("  124020369     ", {"color": MUTED}),
            ("Bihua Wen", {"bold": True}), ("  124090670", {"color": MUTED}),
        ],
        0.70, 3.04, 6.8, 0.30, size=12.8,
    )
    add_rich_line(slide, [("Mingrui Liu", {"bold": True}), ("  124090375", {"color": MUTED})],
                  0.70, 3.39, 6.8, 0.30, size=12.8)
    add_rect(slide, 0.68, 4.00, 6.87, 1.87, fill=BLUE_SOFT, line=BLUE_SOFT)
    add_text(slide, "Core claim", 0.96, 4.26, 1.35, 0.30, size=15.0, color=BLUE, bold=True)
    add_text(slide,
             "Separating physical-border motion from displayed-content motion improves both geometric accuracy and temporal stability.",
             0.96, 4.72, 6.27, 0.72, size=17.0, bold=True, line_spacing=1.10)
    add_rect(slide, 7.93, 1.33, 4.78, 2.12, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_01_pipeline.png", 8.05, 1.45, 4.54, 1.84)
    add_rect(slide, 7.93, 3.73, 4.78, 2.14, fill=INK, line=INK)
    add_picture_crop(slide, IMG / "annotated_dataset_mosaic.jpg", 8.05, 3.85, 4.54, 1.90,
                     crop_top=0.00, crop_bottom=0.48)
    add_footer(slide, 1, "paper Figure 1; dataset mosaic")
    add_notes(slide,
              "[About 15 seconds] This project addresses the geometric front end of captured-screen video. Our core claim is simple: the physical screen boundary should drive the homography, while internal content motion should remain independent.")

    # 2 — Problem and goal
    slide = prs.slides.add_slide(blank)
    add_header(slide, "PROBLEM", "A Handheld Camera and the Display Do Not Share One Motion", 2,
               "paper Introduction", title_size=29.5)
    add_rect(slide, 0.66, 1.55, 5.83, 3.52, fill=RED_SOFT, line=RED)
    add_text(slide, "THE CAMERA VIEW ADDS", 0.96, 1.86, 3.5, 0.32,
             size=15.0, color=RED, bold=True)
    add_bullet_list(slide, [
        "Perspective distortion and off-screen background",
        "Handheld translation, rotation, and scale change",
        "Glare, weak borders, partial visibility, and clutter",
    ], 0.96, 2.35, 5.10, 1.96, size=15.5, bullet_color=RED, gap=0.68)

    add_rect(slide, 6.78, 1.55, 5.90, 3.52, fill=BLUE_SOFT, line=BLUE)
    add_text(slide, "THE DISPLAY CONTENT ADDS", 7.08, 1.86, 3.8, 0.32,
             size=15.0, color=BLUE, bold=True)
    add_bullet_list(slide, [
        "Page scrolling and cursor motion",
        "Video playback and animated regions",
        "Coherent texture motion unrelated to the physical screen",
    ], 7.08, 2.35, 5.20, 1.96, size=15.5, bullet_color=BLUE, gap=0.68)

    add_rect(slide, 0.66, 5.34, 12.02, 1.26, fill=BLUE_PALE, line=LINE)
    add_rich_line(slide, [("Goal. ", {"color": BLUE, "bold": True}),
                          ("Continuously estimate the physical screen plane and render a stable frontal video while preserving true content motion.", {"bold": True})],
                  0.96, 5.68, 11.42, 0.58, size=17.0, valign=MSO_ANCHOR.MIDDLE)
    add_notes(slide,
              "[About 20 seconds] A handheld recording mixes two different sources of motion. The camera changes the physical screen geometry, while the displayed page or video can move independently. The task is to remove camera-induced plane changes without suppressing genuine content motion.")

    # 3 — Scope and contributions
    slide = prs.slides.add_slide(blank)
    add_header(slide, "SCOPE & CONTRIBUTIONS", "A Focused Geometry Stage Before Restoration", 3,
               "paper Introduction and Related Work", title_size=30)
    stages = [
        (0.66, "RAW CAPTURE", "Background + perspective + hand motion", RED_SOFT, RED),
        (4.55, "THIS PROJECT", "Locate, rectify, and stabilize the physical screen", BLUE_SOFT, BLUE),
        (8.44, "DOWNSTREAM", "OCR · restoration · demoiréing · archival", TEAL_SOFT, TEAL),
    ]
    for i, (x, label, body, fill, accent) in enumerate(stages):
        add_rect(slide, x, 1.58, 3.55, 1.30, fill=fill, line=accent)
        add_text(slide, label, x + 0.18, 1.78, 3.19, 0.25, size=12.5, color=accent,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.22, 2.17, 3.11, 0.46, size=14.0, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < 2:
            add_arrow(slide, x + 3.62, 2.23, x + 3.86, 2.23, color=LINE, width=2.2)

    add_text(slide, "Three contributions", 0.68, 3.31, 3.2, 0.34,
             size=16.0, color=BLUE, bold=True)
    contributions = [
        (0.66, "01", "Border-driven estimation", "A complete captured-screen pipeline in which physical boundary evidence drives the screen-plane trajectory."),
        (4.55, "02", "Reproducible evaluation", "Corner annotations, shared metric code, and five capture conditions with common initialization and output settings."),
        (8.44, "03", "Evidence across failure modes", "Geometry and temporal improvements are largest on scrolling pages and weak-border scenes."),
    ]
    for x, idx, title, body in contributions:
        add_rect(slide, x, 3.78, 3.55, 2.32, fill=BLUE_PALE, line=LINE)
        add_text(slide, idx, x + 0.20, 4.02, 0.45, 0.26, size=11.5, color=BLUE, bold=True)
        add_text(slide, title, x + 0.66, 3.98, 2.65, 0.38, size=15.2, bold=True)
        add_text(slide, body, x + 0.20, 4.56, 3.10, 1.15, size=13.2, color=MUTED,
                 line_spacing=1.10)
    add_notes(slide,
              "[About 20 seconds] The project stops at geometry. It produces the stable screen region required by later OCR, restoration, or demoiréing. The paper contributes a border-driven pipeline, a reproducible five-condition evaluation, and evidence that the design is especially useful under scrolling and weak boundaries.")

    # 4 — Why border evidence
    slide = prs.slides.add_slide(blank)
    add_header(slide, "KEY INSIGHT", "Three Methods Trust Three Different Sources of Evidence", 4,
               "paper Sections 1 and 3.4; Figure 5", title_size=29)
    cards = [
        (0.66, "FRAME-WISE", "Trusts the current frame", "Independent detection\nNoise becomes jitter", RED, RED_SOFT),
        (4.55, "OPTICAL FLOW", "Trusts internal texture", "Content motion propagates\nScrolling becomes drift", ORANGE, "FFF3E8"),
        (8.44, "BORDER-GUIDED", "Trusts the physical boundary", "Border defines the plane\nFlow diagnoses conflict", BLUE, BLUE_SOFT),
    ]
    for x, label, belief, consequence, accent, fill in cards:
        add_rect(slide, x, 1.55, 3.55, 2.10, fill=fill, line=accent)
        add_text(slide, label, x + 0.22, 1.80, 3.11, 0.28, size=15.2, color=accent, bold=True)
        add_text(slide, belief, x + 0.22, 2.24, 3.11, 0.28, size=14.0, color=MUTED, bold=True)
        add_line(slide, x + 0.22, 2.66, x + 3.23, 2.66, color=LINE, width=0.8)
        add_text(slide, consequence, x + 0.22, 2.84, 3.11, 0.58, size=14.2, bold=True,
                 line_spacing=1.12)

    add_rect(slide, 0.66, 3.96, 12.02, 2.47, fill=WHITE, line=LINE)
    add_picture_crop(slide, IMG / "figure_05_qualitative.png", 0.80, 4.10, 11.74, 2.03,
                     crop_top=0.205, crop_bottom=0.615)
    add_text(slide, "Same scrolling moment: input annotation · frame-wise · optical flow · proposed",
             0.90, 6.16, 11.50, 0.24, size=11.5, color=MUTED, italic=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 30 seconds] The comparison methods differ mainly in what they trust. Frame-wise detection trusts each image independently. Optical flow trusts interior texture over time. Our method trusts the physical screen border and uses interior flow only to flag a content-motion conflict. This distinction is most visible on a scrolling page.")

    # 5 — Pipeline overview
    slide = prs.slides.add_slide(blank)
    add_header(slide, "METHOD OVERVIEW", "Border Evidence Drives the Homography", 5,
               "paper Figure 1", title_size=31)
    add_rect(slide, 0.62, 1.52, 12.10, 4.42, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_01_pipeline.png", 0.78, 1.67, 11.78, 4.12)
    add_rect(slide, 0.66, 6.14, 12.02, 0.63, fill=BLUE_SOFT, line=BLUE_SOFT)
    add_text(slide,
             "Previous quadrilateral → local border evidence → four fitted lines → valid quadrilateral → diagnostics → smoothing and warp",
             0.88, 6.29, 11.58, 0.32, size=15.0, color=BLUE_DARK, bold=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 20 seconds] Starting from initial corners, the previous quadrilateral predicts local search bands. Physical edge evidence produces four fitted lines and their intersections. LK consistency, gating, fallback, smoothing, and warping complete the pipeline, but the boundary remains the main geometry cue.")

    # 6 — Frame update and safety logic
    slide = prs.slides.add_slide(blank)
    add_header(slide, "FRAME UPDATE", "Boundary Observation Plus Explicit Safety Logic", 6,
               "paper Sections 3.2–3.3", title_size=29.5)
    add_rect(slide, 0.66, 1.55, 6.20, 4.95, fill=BLUE_PALE, line=LINE)
    add_text(slide, "A. RECOVER FOUR PHYSICAL SIDES", 0.94, 1.82, 3.90, 0.28,
             size=14.0, color=BLUE, bold=True)
    # stylized frame / quadrilateral
    add_rect(slide, 1.18, 2.42, 3.20, 2.08, fill=WHITE, line=LINE, radius=False)
    pts = [(1.52, 2.70), (4.00, 2.58), (4.16, 4.14), (1.38, 4.28)]
    for i in range(4):
        x1, y1 = pts[i]
        x2, y2 = pts[(i + 1) % 4]
        add_line(slide, x1, y1, x2, y2, color=BLUE, width=2.8)
        add_dot(slide, x1 - 0.06, y1 - 0.06, 0.12, BLUE)
    # search bands
    add_rect(slide, 1.34, 2.48, 2.75, 0.18, fill=TEAL_SOFT, line=TEAL, radius=False,
             line_width=0.5, transparency=18)
    add_rect(slide, 1.30, 4.16, 2.92, 0.18, fill=TEAL_SOFT, line=TEAL, radius=False,
             line_width=0.5, transparency=18)
    add_text(slide, "predicted search bands", 1.38, 4.58, 2.80, 0.24,
             size=11.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, [
        "Sample profiles along inward normals",
        "Select strong gradients near the predicted side",
        "Robustly fit four lines and intersect neighbors",
    ], 4.63, 2.40, 1.92, 2.25, size=13.0, gap=0.72)

    add_rect(slide, 7.12, 1.55, 5.56, 4.95, fill=WHITE, line=LINE)
    add_text(slide, "B. ACCEPT, RECOVER, OR HOLD", 7.42, 1.82, 4.20, 0.28,
             size=14.0, color=BLUE, bold=True)
    flow = [
        (7.50, 2.32, 4.80, 0.62, "Four usable sides?", BLUE_SOFT, BLUE),
        (7.50, 3.16, 4.80, 0.62, "Valid, convex, geometrically plausible?", BLUE_PALE, LINE),
        (7.50, 4.00, 4.80, 0.62, "Accept boundary + record LK conflict", TEAL_SOFT, TEAL),
        (7.50, 4.84, 2.20, 0.72, "If missing:\nautomatic redetection", RED_SOFT, RED),
        (10.10, 4.84, 2.20, 0.72, "If still missing:\ncarry forward last valid", BLUE_PALE, LINE),
    ]
    for x, y, w, h, txt, fill, accent in flow:
        add_rect(slide, x, y, w, h, fill=fill, line=accent)
        add_text(slide, txt, x + 0.14, y + 0.10, w - 0.28, h - 0.16,
                 size=13.0, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_arrow(slide, 9.90, 2.95, 9.90, 3.12, color=LINE, width=1.8)
    add_arrow(slide, 9.90, 3.79, 9.90, 3.96, color=LINE, width=1.8)
    add_notes(slide,
              "[About 30 seconds] Each side is recovered locally around its predicted position. Profiles are sampled along inward normals, strong nearby gradients are selected, and four robust lines are intersected. A candidate is accepted only when all sides are usable and the quadrilateral is valid and plausible. Missing evidence triggers redetection; only a second failure carries forward the last valid plane.")

    # 7 — Dataset and protocol
    slide = prs.slides.add_slide(blank)
    add_header(slide, "DATASET & PROTOCOL", "Five Capture Conditions, One Shared Evaluation", 7,
               "paper Section 4.1; dataset mosaic", title_size=30)
    add_metric_card(slide, 0.66, 1.53, 2.04, 1.24, "50", "self-collected clips", BLUE, BLUE_SOFT, 27, 12.5)
    add_metric_card(slide, 2.90, 1.53, 2.32, 1.24, "14,985", "total frames", BLUE, BLUE_SOFT, 25, 12.5)
    add_metric_card(slide, 5.42, 1.53, 2.04, 1.24, "10", "annotated evaluation clips", BLUE, BLUE_SOFT, 27, 11.5)

    add_text(slide, "Capture conditions", 0.68, 3.03, 2.40, 0.32,
             size=16.0, color=BLUE, bold=True)
    classes = ["Static", "Scrolling", "Screen video", "Weak border", "Challenging"]
    for i, label in enumerate(classes):
        x = 0.66 + i * 1.37
        fill = RED_SOFT if i == 4 else BLUE_PALE
        accent = RED if i == 4 else LINE
        add_rect(slide, x, 3.48, 1.20, 1.08, fill=fill, line=accent)
        add_text(slide, label, x + 0.08, 3.72, 1.04, 0.50,
                 size=12.2, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 0.66, 4.85, 6.80, 1.48, fill=BLUE_PALE, line=LINE)
    add_text(slide, "Shared protocol", 0.94, 5.10, 1.75, 0.28,
             size=14.5, color=BLUE, bold=True)
    add_text(slide,
             "Same input · initialization · output canvas · annotations · encoder · metric code\nFrame 0 supplies initialization and is excluded from geometry scoring.",
             0.94, 5.49, 6.18, 0.60, size=13.2, bold=True, line_spacing=1.12)

    add_rect(slide, 7.76, 1.53, 4.92, 4.80, fill=INK, line=INK)
    add_picture_contain(slide, IMG / "annotated_dataset_mosaic.jpg", 7.87, 1.64, 4.70, 4.56)
    add_notes(slide,
              "[About 20 seconds] The collection contains fifty clips and 14,985 frames across five conditions. Ten clips, two per condition, form the annotated evaluation split. Every method shares the same initialization, output canvas, annotations, encoder, and metric implementation.")

    # 8 — Metrics
    slide = prs.slides.add_slide(blank)
    add_header(slide, "METRICS", "Accuracy, Stability, and Signal Preservation", 8,
               "paper Section 4.2", title_size=31)
    metric_groups = [
        (0.66, "GEOMETRY", BLUE, BLUE_SOFT,
         ["Corner RMSE ↓", "Quadrilateral IoU ↑", "Aspect-ratio error ↓"],
         "Is the estimated screen in the correct position and shape?"),
        (4.55, "TRAJECTORY", TEAL, TEAL_SOFT,
         ["Translation variation ↓", "Rotation variation ↓", "Scale variation ↓"],
         "Does the estimated screen plane change smoothly over time?"),
        (8.44, "SIGNAL DIAGNOSTICS", ORANGE, "FFF3E8",
         ["SSIM · gradient similarity · Edge F1", "FFT and orientation similarity", "High-frequency / band energy ratios"],
         "Does rectification preserve the captured screen signal?"),
    ]
    for x, label, accent, fill, metrics, question in metric_groups:
        add_rect(slide, x, 1.58, 3.55, 4.52, fill=fill, line=accent)
        add_text(slide, label, x + 0.24, 1.88, 3.07, 0.34,
                 size=15.0, color=accent, bold=True, align=PP_ALIGN.CENTER)
        add_line(slide, x + 0.28, 2.42, x + 3.27, 2.42, color=LINE, width=0.8)
        add_bullet_list(slide, metrics, x + 0.30, 2.72, 2.98, 1.92,
                        size=15.0, bullet_color=accent, gap=0.62)
        add_text(slide, question, x + 0.28, 4.90, 2.99, 0.78,
                 size=13.0, color=MUTED, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0.66, 6.30, 12.02, 0.47, fill=BLUE_PALE, line=BLUE_PALE)
    add_text(slide,
             "Signal metrics are reference-based diagnostics on annotated frames; they do not measure demoiréing and are not primary ranking metrics.",
             0.86, 6.40, 11.62, 0.26, size=12.4, color=MUTED, italic=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 20 seconds] Geometry metrics test position and shape, trajectory metrics test temporal change, and signal diagnostics test whether the warped result preserves the captured screen structure. The signal panel is explicitly diagnostic: it does not claim demoiréing performance.")

    # 9 — Overall results
    slide = prs.slides.add_slide(blank)
    add_header(slide, "OVERALL RESULTS", "The Border-Guided Method Is More Accurate and More Stable", 9,
               "paper Figure 2", title_size=28.5)
    add_rect(slide, 0.64, 1.52, 9.26, 5.03, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_02_overall_results.png", 0.78, 1.67, 8.98, 4.72)
    add_metric_card(slide, 10.16, 1.55, 2.52, 1.42, "3.87 px", "Corner RMSE ↓", BLUE, BLUE_SOFT, 25, 13)
    add_metric_card(slide, 10.16, 3.15, 2.52, 1.42, "0.996", "Quadrilateral IoU ↑", TEAL, TEAL_SOFT, 25, 12.5)
    add_metric_card(slide, 10.16, 4.75, 2.52, 1.42, "2.45", "px/frame translation ↓", BLUE, BLUE_SOFT, 25, 11.5)
    add_notes(slide,
              "[About 30 seconds] Across ten clip-level summaries, the proposed method reaches 3.87-pixel corner RMSE, 0.996 IoU, and 2.45 pixels per frame of translation variation. Geometry and temporal stability improve together, so the result is not created by simply freezing the trajectory.")

    # 10 — Category results
    slide = prs.slides.add_slide(blank)
    add_header(slide, "RESULTS BY CONDITION", "The Largest Gains Occur Under Scrolling and Weak Borders", 10,
               "paper Figure 3", title_size=28.5)
    add_rect(slide, 0.64, 1.52, 9.20, 5.03, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_03_category_results.png", 0.78, 1.66, 8.92, 4.74)
    add_metric_card(slide, 10.08, 1.56, 2.60, 1.50, "2.87 px", "Scrolling RMSE\n31.76 frame-wise · 81.67 flow", BLUE, BLUE_SOFT, 24, 11.5)
    add_metric_card(slide, 10.08, 3.24, 2.60, 1.50, "9.35 px", "Weak-border RMSE\nBoth comparisons >155 px", TEAL, TEAL_SOFT, 24, 11.5)
    add_rect(slide, 10.08, 4.90, 2.60, 1.55, fill=RED_SOFT, line=RED)
    add_text(slide, "LIMIT", 10.30, 5.10, 2.16, 0.25, size=12.0, color=RED, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Challenging scenes: geometry is slightly worse than frame-wise detection, but the trajectory is more stable.",
             10.30, 5.44, 2.16, 0.82, size=11.0, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 30 seconds] Scrolling creates the strongest content-motion conflict: optical-flow RMSE rises to 81.67 pixels, while the proposed method reaches 2.87. Weak borders also show a large gain, with 9.35 pixels versus more than 155 for both comparison methods. Challenging scenes remain the main limitation.")

    # 11 — Per-clip consistency
    slide = prs.slides.add_slide(blank)
    add_header(slide, "PER-CLIP CONSISTENCY", "The Method Does Not Rely on Freezing the Trajectory", 11,
               "paper Figure 4", title_size=29.5)
    add_rect(slide, 0.64, 1.52, 12.04, 4.63, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_04_proposed_clip_results.png", 0.80, 1.65, 11.72, 4.34)
    chips = [
        (0.72, "7 / 10", "clips below 5 px RMSE", BLUE),
        (4.56, "10 / 10", "clips below 15 px RMSE", TEAL),
        (8.40, "0", "held frames in the evaluated run", BLUE),
    ]
    for x, value, label, accent in chips:
        add_rect(slide, x, 6.27, 3.55, 0.50, fill=BLUE_PALE, line=accent)
        add_rich_line(slide, [(value, {"color": accent, "bold": True, "size": 17}),
                              (f"   {label}", {"bold": True, "size": 13})],
                      x + 0.20, 6.36, 3.15, 0.28, size=13, align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 25 seconds] Seven of ten clips remain below five-pixel RMSE and all ten remain below fifteen. The evaluated run contains zero held frames, so the stability is not produced by refusing updates. The higher-error clips are concentrated in weak-border and challenging conditions.")

    # 12 — Qualitative comparison
    slide = prs.slides.add_slide(blank)
    add_header(slide, "QUALITATIVE COMPARISON", "The Screen Extent Is Preserved Under Content Motion", 12,
               "paper Figure 5", title_size=29.5)
    add_text(slide, "Scrolling page", 0.72, 1.49, 2.20, 0.28,
             size=15.5, color=BLUE, bold=True)
    add_rect(slide, 0.66, 1.84, 12.02, 2.14, fill=WHITE, line=LINE)
    add_picture_crop(slide, IMG / "figure_05_qualitative.png", 0.82, 1.98, 11.70, 1.84,
                     crop_top=0.205, crop_bottom=0.615)
    add_text(slide, "Weak-border example", 0.72, 4.21, 2.50, 0.28,
             size=15.5, color=TEAL, bold=True)
    add_rect(slide, 0.66, 4.56, 12.02, 2.08, fill=WHITE, line=LINE)
    add_picture_crop(slide, IMG / "figure_05_qualitative.png", 0.82, 4.70, 11.70, 1.76,
                     crop_top=0.635, crop_bottom=0.165)
    add_notes(slide,
              "[About 25 seconds] Qualitative outputs match the metrics. On scrolling content, the border-guided result preserves the physical screen instead of following the page. In the weak-border example, it also maintains the full screen extent where independent detection can shift the crop.")

    # 13 — Signal diagnostics
    slide = prs.slides.add_slide(blank)
    add_header(slide, "SIGNAL PRESERVATION", "Rectification Preserves Local and Frequency Structure", 13,
               "paper Figure 6", title_size=29)
    add_rect(slide, 0.64, 1.52, 9.42, 5.06, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_06_signal_preservation.png", 0.78, 1.66, 9.14, 4.78)
    add_metric_card(slide, 10.30, 1.55, 2.38, 1.24, "0.890", "SSIM", TEAL, TEAL_SOFT, 24, 13)
    add_metric_card(slide, 10.30, 2.96, 2.38, 1.24, "0.930", "Gradient-map similarity", TEAL, TEAL_SOFT, 24, 11.5)
    add_metric_card(slide, 10.30, 4.37, 2.38, 1.24, "0.952", "Edge F1", TEAL, TEAL_SOFT, 24, 13)
    add_rect(slide, 10.30, 5.80, 2.38, 0.76, fill=BLUE_PALE, line=LINE)
    add_text(slide, "Diagnostic only—not a demoiréing claim", 10.48, 6.00, 2.02, 0.34,
             size=11.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 25 seconds] On the representative scrolling clip, the proposed output has higher local-structure similarity: SSIM 0.890, gradient similarity 0.930, and Edge F1 0.952. Frequency similarity and energy ratios also remain close to the annotation-warped reference. These are preservation diagnostics, not demoiréing results.")

    # 14 — Ablation
    slide = prs.slides.add_slide(blank)
    add_header(slide, "ABLATION", "Physical Borders Are Decisive; Smoothing Adds Temporal Stability", 14,
               "paper Table 4; proposal_border_ablation_2026-07-14.md", title_size=27.5)
    table_shape = slide.shapes.add_table(6, 4, Inches(0.66), Inches(1.60), Inches(8.90), Inches(4.80))
    table = table_shape.table
    for idx, width in enumerate([3.65, 1.60, 1.60, 2.05]):
        table.columns[idx].width = Inches(width)
    for i in range(6):
        table.rows[i].height = Inches(0.80)
    for j, header in enumerate(["Variant", "RMSE ↓", "IoU ↑", "Translation ↓"]):
        style_cell(table.cell(0, j), header, BLUE_DARK, color=WHITE, bold=True, size=12.8)
    rows = [
        ("Full Profile-border method", "3.253", "0.996038", "0.752", BLUE_SOFT),
        ("No trajectory smoothing", "2.932", "0.996585", "1.430", RED_SOFT),
        ("No physical border: optical flow", "76.114", "0.916022", "2.205", "FFF3E8"),
        ("LSD border observation", "3.604", "0.995716", "0.961", "F6F8FB"),
        ("Hough border observation", "27.335", "0.974200", "0.897", "F6F8FB"),
    ]
    for i, (name, rmse, iou, trans, fill) in enumerate(rows, 1):
        for j, value in enumerate([name, rmse, iou, trans]):
            style_cell(table.cell(i, j), value, fill,
                       color=(BLUE if i == 1 and j > 0 else INK),
                       bold=(i in (1, 2, 3)), size=(12.4 if j == 0 else 12.8),
                       align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))

    add_rect(slide, 9.84, 1.60, 2.84, 1.40, fill=BLUE_SOFT, line=BLUE)
    add_text(slide, "0.752 → 1.430", 10.04, 1.87, 2.44, 0.38,
             size=21.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "No smoothing nearly doubles translation variation", 10.04, 2.37, 2.44, 0.38,
             size=11.8, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.84, 3.22, 2.84, 1.40, fill=RED_SOFT, line=RED)
    add_text(slide, "3.253 → 76.114", 10.04, 3.49, 2.44, 0.38,
             size=20.0, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Removing physical-border evidence causes content-driven failure", 10.04, 3.99, 2.44, 0.42,
             size=11.5, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.84, 4.84, 2.84, 1.56, fill=TEAL_SOFT, line=TEAL)
    add_text(slide, "Profile is the default", 10.04, 5.12, 2.44, 0.32,
             size=14.0, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "LSD is close but slower; Hough degrades geometry on this clip.", 10.04, 5.56, 2.44, 0.54,
             size=11.6, bold=True, align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 30 seconds] The ablation separates the roles of the modules. Removing smoothing nearly doubles translation variation, so filtering mainly adds stability. Removing physical-border evidence raises RMSE from 3.253 to 76.114 because optical flow follows scrolling content. Profile observations are the default: LSD is close but slower, while Hough clearly degrades geometry on this clip.")

    # 15 — Limitations and conclusion
    slide = prs.slides.add_slide(blank)
    add_header(slide, "LIMITATIONS & CONCLUSION", "The Boundary Is the Strength—and the Main Limitation", 15,
               "paper Discussion and Conclusion", title_size=29)
    add_rect(slide, 0.66, 1.56, 5.85, 2.20, fill=RED_SOFT, line=RED)
    add_text(slide, "LIMITATIONS", 0.96, 1.86, 2.0, 0.32, size=15.0, color=RED, bold=True)
    add_bullet_list(slide, [
        "Very dark or low-contrast borders",
        "Reflections, glare, and partial occlusion",
        "Sparse keyframe annotations rather than dense ground truth",
    ], 0.96, 2.34, 5.10, 1.20, size=14.4, bullet_color=RED, gap=0.43)

    add_rect(slide, 6.78, 1.56, 5.90, 2.20, fill=BLUE_SOFT, line=BLUE)
    add_text(slide, "FINAL CONCLUSION", 7.08, 1.86, 2.4, 0.32, size=15.0, color=BLUE, bold=True)
    add_text(slide,
             "Physical screen edges are a better primary cue for the screen plane than independent frame detection or internal-content optical flow.",
             7.10, 2.34, 5.24, 1.04, size=17.0, bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE, line_spacing=1.08)

    add_text(slide, "NEXT STEPS", 0.68, 4.13, 1.65, 0.30, size=15.0, color=BLUE, bold=True)
    next_steps = [
        (0.66, "Multi-cue boundary models", "Fuse profiles, long lines, color differences, and rectangular constraints."),
        (4.55, "Per-edge confidence", "Estimate four independent confidence values and complete weak sides across frames."),
        (8.44, "Active recovery", "Reinitialize after repeated failures and expand glare/occlusion annotations."),
    ]
    for x, title, body in next_steps:
        add_rect(slide, x, 4.58, 3.55, 1.42, fill=BLUE_PALE, line=LINE)
        add_text(slide, title, x + 0.20, 4.82, 3.15, 0.32, size=15.0, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.20, 5.30, 3.15, 0.48, size=12.8, color=MUTED,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "3.87 px RMSE   ·   0.996 IoU   ·   2.45 px/frame",
             0.68, 6.26, 11.98, 0.42, size=19.0, color=BLUE, bold=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "[About 20 seconds] The method remains limited by weak or occluded boundaries and by sparse annotations. The direct next steps are multi-cue boundary models, per-edge confidence, and stronger recovery. The final result is 3.87-pixel RMSE, 0.996 IoU, and 2.45 pixels per frame of translation variation.")

    return prs


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT)
    print(f"Saved {len(prs.slides)} slides to {OUT}")


if __name__ == "__main__":
    main()
