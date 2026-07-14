from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
IMG = ROOT / "images"
OUT = ROOT / "exports" / "screen_normalize_final_6min.pptx"

W = 13.333
H = 7.5

BLUE = "1F57C4"
BLUE_DARK = "17479E"
BLUE_SOFT = "E8EFFB"
BLUE_PALE = "F4F6FA"
INK = "1A1A1A"
MUTED = "555B66"
LINE = "D3DAE6"
RED = "C43D2E"
RED_SOFT = "FCEDEA"
ORANGE = "E97822"
TEAL = "2E8579"
TEAL_SOFT = "E6F3F1"
WHITE = "FFFFFF"

FONT_CN = "Microsoft YaHei"
FONT_LATIN = "Aptos"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def set_run_font(run, size: float, color: str = INK, bold: bool = False,
                 italic: bool = False, font: str = FONT_CN) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.font.bold = bold
    run.font.italic = italic

    r_pr = run._r.get_or_add_rPr()
    ea = r_pr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        r_pr.append(ea)
    ea.set("typeface", FONT_CN)


def set_text_margins(shape, left=0.08, right=0.08, top=0.04, bottom=0.04):
    tf = shape.text_frame
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)


def add_text(slide, text: str, x: float, y: float, w: float, h: float,
             size: float = 16, color: str = INK, bold: bool = False,
             italic: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, font: str = FONT_CN,
             margin: float = 0.0, line_spacing: float = 1.08):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text_margins(shape, margin, margin, margin, margin)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        r = p.add_run()
        r.text = line
        set_run_font(r, size, color, bold, italic, font)
    return shape


def add_rich_line(slide, parts, x, y, w, h, size=16, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, margin=0.0):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text_margins(shape, margin, margin, margin, margin)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    for part in parts:
        if isinstance(part, str):
            text = part
            kwargs = {}
        else:
            text = part[0]
            kwargs = part[1] if len(part) > 1 and isinstance(part[1], dict) else {}
        r = p.add_run()
        r.text = text
        set_run_font(
            r,
            kwargs.get("size", size),
            kwargs.get("color", INK),
            kwargs.get("bold", False),
            kwargs.get("italic", False),
            kwargs.get("font", FONT_CN),
        )
    return shape


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True,
             line_width=0.8, transparency=0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_line(slide, x1, y1, x2, y2, color=BLUE, width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=2.0):
    line = add_line(slide, x1, y1, x2, y2, color, width)
    line.line.end_arrowhead = True
    return line


def add_circle(slide, cx, cy, d, fill=BLUE, line=BLUE):
    return add_rect(slide, cx - d / 2, cy - d / 2, d, d, fill, line, radius=True, line_width=0.8)


def image_size(path: Path):
    with Image.open(path) as im:
        return im.size


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float,
                        border: str | None = None, border_width: float = 0.8):
    px_w, px_h = image_size(path)
    scale = min(w / px_w, h / px_h)
    draw_w = px_w * scale
    draw_h = px_h * scale
    pic = slide.shapes.add_picture(
        str(path),
        Inches(x + (w - draw_w) / 2),
        Inches(y + (h - draw_h) / 2),
        width=Inches(draw_w),
        height=Inches(draw_h),
    )
    if border:
        pic.line.color.rgb = rgb(border)
        pic.line.width = Pt(border_width)
    return pic


def add_picture_crop(slide, path: Path, x: float, y: float, w: float, h: float,
                     crop_left=0.0, crop_top=0.0, crop_right=0.0, crop_bottom=0.0,
                     border: str | None = None):
    pic = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )
    pic.crop_left = crop_left
    pic.crop_top = crop_top
    pic.crop_right = crop_right
    pic.crop_bottom = crop_bottom
    if border:
        pic.line.color.rgb = rgb(border)
        pic.line.width = Pt(0.8)
    return pic


def add_footer(slide, page: int, source: str = ""):
    add_line(slide, 0.62, 7.02, 12.70, 7.02, color=LINE, width=0.5)
    add_text(slide, "ECE4512 Final Project · 2026", 0.62, 7.08, 1.75, 0.19,
             size=7.2, color=MUTED)
    if source:
        add_text(slide, f"来源：{source}", 2.42, 7.08, 9.68, 0.19,
                 size=6.7, color=MUTED)
    add_text(slide, str(page), 12.45, 7.08, 0.25, 0.19, size=7.6, color=MUTED,
             align=PP_ALIGN.RIGHT)


def add_header(slide, kicker: str, title: str, page: int, source: str = "",
               title_size: float = 29):
    add_text(slide, kicker.upper(), 0.62, 0.32, 8.6, 0.28, size=11.5, color=BLUE, bold=True)
    add_text(slide, title, 0.62, 0.67, 12.0, 0.58, size=title_size, bold=True,
             valign=MSO_ANCHOR.MIDDLE)
    add_line(slide, 0.64, 1.26, 1.55, 1.26, color=BLUE, width=2.4)
    add_footer(slide, page, source)


def add_notes(slide, text: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def style_cell(cell, text, fill, color=INK, bold=False, size=11.5,
               align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill)
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run_font(r, size, color, bold)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # Slide 1 — Cover
    slide = prs.slides.add_slide(blank)
    add_text(slide, "ECE4512 FINAL PROJECT · 2026", 0.70, 0.48, 6.7, 0.30,
             size=12.2, color=BLUE, bold=True)
    add_text(slide, "物理边框主导的\n拍屏视频几何归一化", 0.70, 0.92, 6.95, 1.35,
             size=28.5, bold=True, line_spacing=1.03)
    add_text(slide, "Border-guided screen-plane normalization", 0.72, 2.36, 6.5, 0.34,
             size=14.2, color=MUTED, italic=True, font=FONT_LATIN)
    add_rich_line(
        slide,
        [
            ("Rongshuo Wen", {"bold": True}), ("  124020369      ", {"color": MUTED}),
            ("Bihua Wen", {"bold": True}), ("  124090670", {"color": MUTED}),
        ],
        0.72, 2.85, 6.6, 0.30, size=11.8,
    )
    add_rich_line(
        slide,
        [("Mingrui Liu", {"bold": True}), ("  124090375", {"color": MUTED})],
        0.72, 3.18, 6.6, 0.30, size=11.8,
    )

    add_rect(slide, 0.70, 3.72, 6.85, 2.35, fill=BLUE_SOFT, line=BLUE_SOFT)
    add_rich_line(slide, [("核心问题。", {"color": BLUE, "bold": True}),
                          ("拍屏视频中的页面滚动、屏内视频与相机运动并不相同。")],
                  0.96, 4.00, 6.25, 0.42, size=15.5)
    add_text(slide,
             "我们让物理屏幕边框直接决定单应矩阵，内部 LK/RANSAC 仅用于一致性诊断，"
             "从而在保留真实内容运动的同时，输出稳定的正面屏幕视频。",
             0.96, 4.53, 6.18, 1.16, size=15.2, color=INK, line_spacing=1.16)

    add_rect(slide, 8.04, 1.40, 4.67, 1.88, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_01_pipeline.png", 8.12, 1.49, 4.51, 1.69)
    add_text(slide, "项目现有方法图：边框驱动的屏幕平面归一化", 8.15, 3.08, 4.45, 0.20,
             size=8.2, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 8.04, 3.58, 4.67, 2.48, fill=WHITE, line=LINE)
    add_picture_crop(slide, IMG / "annotated_dataset_mosaic.jpg", 8.12, 3.66, 4.51, 2.20,
                     crop_top=0.00, crop_bottom=0.48)
    add_text(slide, "自采集拍屏样例与人工四角标注", 8.15, 5.84, 4.45, 0.20,
             size=8.2, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 1, "figure_01_pipeline.png；annotated_dataset_mosaic.jpg")
    add_notes(slide,
              "【约 25 秒】本项目研究的是拍屏视频的几何前端。输入是包含背景、透视畸变、"
              "手持抖动以及屏幕内部动态内容的完整视频；输出是稳定、正面的屏幕视频。"
              "核心思路是让物理屏幕边框决定单应矩阵，而不是让页面滚动或屏内视频带着屏幕平面一起移动。")

    # Slide 2 — Problem and insight
    slide = prs.slides.add_slide(blank)
    add_header(slide, "WHY BORDER EVIDENCE", "问题的本质：内容运动 ≠ 屏幕运动", 2,
               "完善大纲.md；figure_05_qualitative.png")
    cards = [
        (0.65, "逐帧检测", "相信当前帧", "检测噪声 → 抖动\n弱边框 → 定位偏移", RED, RED_SOFT),
        (4.45, "相邻帧光流", "相信内部纹理", "页面滚动 → 内容驱动\n逐帧传播 → 累积漂移", ORANGE, "FFF3E8"),
        (8.25, "边框主导（本文）", "相信物理边界", "边框决定四边形\n内部光流只做诊断", BLUE, BLUE_SOFT),
    ]
    for x, name, belief, weakness, accent, soft in cards:
        add_rect(slide, x, 1.56, 3.46, 2.23, fill=soft, line=accent, line_width=1.1)
        add_text(slide, name, x + 0.22, 1.78, 2.98, 0.36, size=17, color=accent, bold=True)
        add_text(slide, belief, x + 0.22, 2.23, 2.98, 0.30, size=12.5, color=MUTED, bold=True)
        add_line(slide, x + 0.22, 2.64, x + 3.18, 2.64, color=LINE, width=0.8)
        add_text(slide, weakness, x + 0.22, 2.80, 2.98, 0.72, size=13.2, color=INK,
                 line_spacing=1.14)

    add_text(slide, "滚动页面的同一时刻：输入 + 三种方法", 0.70, 4.08, 5.2, 0.28,
             size=11.4, color=BLUE, bold=True)
    add_rect(slide, 0.65, 4.40, 12.03, 2.18, fill=WHITE, line=LINE)
    add_picture_crop(slide, IMG / "figure_05_qualitative.png", 0.76, 4.51, 11.81, 1.92,
                     crop_top=0.205, crop_bottom=0.615)
    add_text(slide, "观察重点：屏幕范围是否完整、边缘是否漂移，以及滚动内容是否被正常保留。",
             0.77, 6.56, 11.78, 0.26, size=10.0, color=MUTED, italic=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "【约 45 秒】问题不只是把画面裁成矩形，而是区分屏幕运动和内容运动。"
              "逐帧检测只相信当前帧，弱边框下容易偏移并产生抖动；相邻帧光流相信内部纹理，"
              "当网页滚动时会把内容运动误当成屏幕运动并累积漂移。我们的方案优先相信物理屏幕边界。"
              "下方是项目现有的滚动页面对比：需要同时检查完整屏幕范围、边缘稳定性和真实滚动是否保留。")

    # Slide 3 — Pipeline
    slide = prs.slides.add_slide(blank)
    add_header(slide, "METHOD", "当前边框主导管线", 3,
               "figure_01_pipeline.png")
    add_rect(slide, 0.62, 1.50, 12.08, 3.63, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_01_pipeline.png", 0.74, 1.60, 11.84, 3.42)

    method_cards = [
        (0.66, "01", "物理边框决定单应矩阵", "在上一帧预测边附近采样 profile，拟合四条边线并求交。"),
        (4.55, "02", "LK/RANSAC 只做诊断", "发现内部运动冲突，但不让滚动内容接管屏幕平面。"),
        (8.44, "03", "门控、重检与平滑", "拒绝非法四边形；必要时重检，并平滑最终轨迹。"),
    ]
    for x, idx, title, body in method_cards:
        add_rect(slide, x, 5.40, 3.60, 1.28, fill=BLUE_PALE, line=LINE)
        add_text(slide, idx, x + 0.18, 5.61, 0.50, 0.30, size=11.2, color=BLUE, bold=True)
        add_text(slide, title, x + 0.66, 5.56, 2.68, 0.34, size=13.1, color=INK, bold=True)
        add_text(slide, body, x + 0.18, 6.02, 3.12, 0.44, size=10.5, color=MUTED,
                 line_spacing=1.08)
    add_notes(slide,
              "【约 55 秒】管线从首帧四角初始化开始。之后利用上一帧四边形预测四条边的局部搜索带，"
              "沿内法线采样梯度 profile，选择靠近预测边的高梯度点并鲁棒拟合四条边线。"
              "相邻边求交得到当前四边形，再通过凸性、边长和帧间变化等门控。"
              "LK/RANSAC 只判断内部纹理是否与边框运动冲突；只要边框有效，仍以物理边界为准。"
              "最后对轨迹插值和平滑，再投影到固定画布。")

    # Slide 4 — Dataset
    slide = prs.slides.add_slide(blank)
    add_header(slide, "EVALUATION", "数据集与实验设置", 4,
               "完善大纲.md；annotated_dataset_mosaic.jpg")

    stat_specs = [
        (0.68, "50", "自采集视频"),
        (2.77, "14,985", "总帧数"),
        (5.16, "10", "带标注评估片段"),
    ]
    for x, number, label in stat_specs:
        add_rect(slide, x, 1.55, 1.87 if x != 2.77 else 2.13, 1.20,
                 fill=BLUE_SOFT, line=BLUE_SOFT)
        add_text(slide, number, x + 0.14, 1.72, (1.59 if x != 2.77 else 1.85), 0.48,
                 size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, label, x + 0.10, 2.24, (1.67 if x != 2.77 else 1.93), 0.24,
                 size=10.2, color=MUTED, align=PP_ALIGN.CENTER)

    add_text(slide, "五类拍摄条件", 0.69, 3.02, 2.2, 0.28, size=13.5, color=BLUE, bold=True)
    classes = ["静态页面", "滚动页面", "屏内视频", "弱边框", "挑战场景"]
    class_colors = [BLUE_SOFT, BLUE_SOFT, BLUE_SOFT, BLUE_SOFT, RED_SOFT]
    class_lines = [LINE, LINE, LINE, LINE, RED]
    for i, label in enumerate(classes):
        x = 0.68 + i * 1.36
        add_rect(slide, x, 3.43, 1.18, 1.22, fill=class_colors[i], line=class_lines[i],
                 line_width=1.0)
        add_text(slide, f"CLASS {i + 1}", x + 0.08, 3.61, 1.02, 0.20,
                 size=7.8, color=(RED if i == 4 else BLUE), bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.08, 4.00, 1.02, 0.30, size=11.2, bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_rect(slide, 0.68, 4.94, 6.58, 1.56, fill=BLUE_PALE, line=LINE)
    add_rich_line(slide, [("公平比较。", {"color": BLUE, "bold": True}),
                          ("三种主方法共享相同输入、初始化、输出画布、标注和指标代码。")],
                  0.92, 5.15, 6.05, 0.35, size=13.1)
    add_text(slide, "几何评价排除用于初始化的第 0 帧；片段内先取中位数，再跨片段汇总。",
             0.92, 5.68, 5.98, 0.51, size=11.5, color=MUTED)

    add_rect(slide, 7.63, 1.55, 5.05, 4.95, fill=INK, line=INK)
    add_picture_contain(slide, IMG / "annotated_dataset_mosaic.jpg", 7.73, 1.65, 4.85, 4.70)
    add_text(slide, "项目内现有数据集马赛克；绿色四边形为人工屏幕边界。",
             7.83, 6.28, 4.65, 0.24, size=8.5, color=LINE, italic=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "【约 40 秒】完整数据集包含五十个自采集拍屏视频，共一万四千九百八十五帧，"
              "五类条件各十个片段。当前正式定量评估使用十个带标注片段，每类两个。"
              "三种方法使用相同的输入、首帧初始化、输出画布、人工角点标注和指标代码；"
              "几何评价排除第零帧，先对每个片段取中位数，再跨片段取中位数。")

    # Slide 5 — Overall results
    slide = prs.slides.add_slide(blank)
    add_header(slide, "MAIN RESULTS", "总体结果：更准，同时更稳", 5,
               "figure_02_overall_results.png")
    add_rect(slide, 0.64, 1.53, 8.42, 4.98, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_02_overall_results.png", 0.75, 1.65, 8.20, 4.72)

    add_rect(slide, 9.32, 1.53, 3.37, 4.98, fill=BLUE_SOFT, line=BLUE_SOFT)
    metrics = [
        ("3.87 px", "角点 RMSE", "最低", BLUE),
        ("0.996", "四边形 IoU", "最高", TEAL),
        ("2.45", "px / frame 平移变化", "最低", BLUE),
    ]
    y = 1.86
    for value, label, tag, c in metrics:
        add_text(slide, value, 9.65, y, 2.72, 0.48, size=23.0, color=c, bold=True)
        add_text(slide, label, 9.65, y + 0.49, 2.20, 0.25, size=10.8, color=INK, bold=True)
        add_rect(slide, 11.77, y + 0.46, 0.57, 0.25, fill=WHITE, line=LINE)
        add_text(slide, tag, 11.79, y + 0.50, 0.51, 0.16, size=7.6, color=c, bold=True,
                 align=PP_ALIGN.CENTER)
        y += 1.08
    add_line(slide, 9.65, 5.06, 12.38, 5.06, color=LINE, width=0.8)
    add_text(slide, "结论", 9.65, 5.30, 0.68, 0.24, size=10.3, color=BLUE, bold=True)
    add_text(slide, "改进不只是轨迹更平滑；更关键的是四边形更贴近人工标注。",
             9.65, 5.63, 2.72, 0.62, size=12.0, color=INK, line_spacing=1.10)
    add_notes(slide,
              "【约 55 秒】这里的三个指标都来自项目现有总体结果图。边框主导方法的角点 RMSE 为 3.87 像素，"
              "而逐帧检测和光流分别约为 30.37 和 31.40。IoU 提升到 0.996，"
              "平移变化降低到每帧 2.45 像素。重要的是，我们不是通过冻结轨迹换取表面稳定："
              "几何精度和时间稳定性同时改善，说明物理边框确实比内部内容更适合作为单应矩阵的主证据。")

    # Slide 6 — Category results
    slide = prs.slides.add_slide(blank)
    add_header(slide, "WHERE IT HELPS", "优势集中在滚动与弱边框场景", 6,
               "figure_03_category_results.png")
    add_rect(slide, 0.64, 1.53, 8.83, 4.98, fill=WHITE, line=LINE)
    add_picture_contain(slide, IMG / "figure_03_category_results.png", 0.76, 1.65, 8.59, 4.72)

    add_rect(slide, 9.72, 1.53, 2.97, 1.58, fill=BLUE_SOFT, line=BLUE)
    add_text(slide, "滚动页面", 9.98, 1.78, 2.40, 0.28, size=14.6, color=BLUE, bold=True)
    add_text(slide, "RMSE  2.87 px", 9.98, 2.18, 2.40, 0.31, size=16.4, color=INK, bold=True)
    add_text(slide, "逐帧 31.76 · 光流 81.67", 9.98, 2.60, 2.40, 0.22,
             size=8.8, color=MUTED)

    add_rect(slide, 9.72, 3.30, 2.97, 1.58, fill=TEAL_SOFT, line=TEAL)
    add_text(slide, "弱边框", 9.98, 3.55, 2.40, 0.28, size=14.6, color=TEAL, bold=True)
    add_text(slide, "RMSE  9.35 px", 9.98, 3.95, 2.40, 0.31, size=16.4, color=INK, bold=True)
    add_text(slide, "另外两种方法均 >155 px", 9.98, 4.37, 2.40, 0.22,
             size=8.8, color=MUTED)

    add_rect(slide, 9.72, 5.07, 2.97, 1.44, fill=BLUE_PALE, line=LINE)
    add_text(slide, "挑战场景仍是限制", 9.98, 5.31, 2.40, 0.26, size=12.6, color=RED, bold=True)
    add_text(slide, "几何略逊于逐帧检测，但平移变化 3.74，仍低于 5.19 / 8.56。",
             9.98, 5.72, 2.40, 0.54, size=10.0, color=MUTED, line_spacing=1.08)
    add_notes(slide,
              "【约 50 秒】分类结果解释了优势来自哪里。滚动页面上，光流会跟随内容移动，"
              "RMSE 上升到 81.67 像素；边框方案保持在物理屏幕上，只有 2.87 像素。"
              "弱边框场景中，上一帧预测的局部搜索带也显著优于全帧独立检测，RMSE 为 9.35，"
              "另外两种方法都超过 155。挑战场景是主要限制：逐帧检测几何误差略低，"
              "但我们的轨迹仍更稳定，因此后续应重点处理反光、遮挡和极低对比度。")

    # Slide 7 — Ablation
    slide = prs.slides.add_slide(blank)
    add_header(slide, "ABLATION", "消融：滤波换来稳定，其余模块提供安全网", 7,
               "proposal_border_ablation_2026-07-14.md")
    rows = 6
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.66), Inches(1.62), Inches(8.25), Inches(4.62))
    table = table_shape.table
    widths = [3.10, 1.54, 1.54, 2.07]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for row_idx in range(rows):
        table.rows[row_idx].height = Inches(0.77)

    headers = ["变体", "RMSE ↓", "IoU ↑", "平移变化 ↓"]
    for j, h_text in enumerate(headers):
        style_cell(table.cell(0, j), h_text, BLUE_DARK, color=WHITE, bold=True, size=11.3)

    data = [
        ("完整方法：Profile 边框", "3.253", "0.996038", "0.752"),
        ("去掉轨迹滤波", "2.932", "0.996585", "1.430"),
        ("去掉 LK 一致性诊断", "3.253", "0.996038", "0.752"),
        ("去掉重新检测回退", "3.253", "0.996038", "0.752"),
        ("放宽边缘门控", "3.253", "0.996038", "0.752"),
    ]
    for i, row in enumerate(data, start=1):
        fill = BLUE_SOFT if i == 1 else (RED_SOFT if i == 2 else "F6F8FB")
        for j, value in enumerate(row):
            style_cell(table.cell(i, j), value, fill,
                       color=(BLUE if i == 1 and j > 0 else INK),
                       bold=(i in (1, 2)), size=(11.2 if j == 0 else 11.0),
                       align=(PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))

    add_rect(slide, 9.20, 1.62, 3.48, 2.05, fill=BLUE_SOFT, line=BLUE)
    add_text(slide, "0.752  →  1.430", 9.48, 1.96, 2.91, 0.44,
             size=22.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "去掉滤波后，平移变化接近翻倍", 9.49, 2.55, 2.89, 0.46,
             size=12.0, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "代价：稀疏标注帧 RMSE 略降，但帧间抖动明显增加。",
             9.48, 3.06, 2.91, 0.37, size=9.2, color=MUTED, align=PP_ALIGN.CENTER)

    add_rect(slide, 9.20, 3.93, 3.48, 2.31, fill=BLUE_PALE, line=LINE)
    add_text(slide, "LK / 重检 / 门控", 9.48, 4.24, 2.91, 0.34,
             size=15.0, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "在该滚动片段上数值不变", 9.48, 4.73, 2.91, 0.28,
             size=11.3, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide,
             "说明 Profile 边框证据逐帧成功；这些模块主要承担异常检测与恢复，而不是正常帧的主估计。",
             9.50, 5.18, 2.87, 0.73, size=10.2, color=MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.10)
    add_text(slide, "表中数值均直接转录自项目现有 Proposal Border Ablation 文档。",
             0.69, 6.47, 8.15, 0.24, size=8.6, color=MUTED, italic=True)
    add_notes(slide,
              "【约 55 秒】消融在代表性的滚动片段上进行。完整方法的 RMSE 为 3.253，"
              "平移变化为 0.752。去掉轨迹滤波后，RMSE 在稀疏标注帧上略降到 2.932，"
              "但平移变化升到 1.430，接近翻倍，所以滤波的作用是时间稳定，而不是提高单帧精度。"
              "去掉 LK 诊断、重检或放宽门控没有改变这个片段的结果，说明 Profile 边框证据每帧都成功；"
              "这些模块是异常情况下的安全网，不会干扰正常主链。")

    # Slide 8 — Limitations and conclusion
    slide = prs.slides.add_slide(blank)
    add_header(slide, "LIMITS & TAKEAWAY", "边界仍是限制，但结论清晰", 8,
               "完善大纲.md §4–5")

    add_rect(slide, 0.66, 1.58, 5.82, 2.03, fill=RED_SOFT, line=RED)
    add_text(slide, "主要局限", 0.94, 1.87, 1.45, 0.31, size=15.0, color=RED, bold=True)
    add_text(slide, "• 弱边缘 / 极低对比度：真实梯度不足，背景纹理可能成为更强候选\n"
                    "• 强反光 / 遮挡：伪梯度会偏移边线，或触发重检与保持",
             0.94, 2.35, 5.16, 0.89, size=12.2, color=INK, line_spacing=1.16)

    add_rect(slide, 6.77, 1.58, 5.91, 2.03, fill=BLUE_SOFT, line=BLUE)
    add_text(slide, "最终结论", 7.07, 1.87, 1.45, 0.31, size=15.0, color=BLUE, bold=True)
    add_text(slide,
             "物理屏幕边框比单帧独立检测或内部内容光流，更适合作为屏幕平面的主要证据。",
             7.07, 2.35, 5.28, 0.86, size=17.4, color=INK, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.10)

    add_text(slide, "下一步", 0.69, 4.00, 1.20, 0.28, size=13.5, color=BLUE, bold=True)
    next_steps = [
        (0.68, "01", "多线索边界融合", "联合 Profile、长线段、颜色差异与矩形约束。"),
        (4.57, "02", "边级置信度", "对四条边分别估计可信度，并跨帧补全低置信度边。"),
        (8.46, "03", "更强失败恢复", "连续失败时主动重新初始化，并扩充反光、遮挡标注。"),
    ]
    for x, idx, title, body in next_steps:
        add_rect(slide, x, 4.42, 3.58, 1.62, fill=BLUE_PALE, line=LINE)
        add_text(slide, idx, x + 0.18, 4.65, 0.48, 0.26, size=10.5, color=BLUE, bold=True)
        add_text(slide, title, x + 0.65, 4.61, 2.65, 0.32, size=13.2, color=INK, bold=True)
        add_text(slide, body, x + 0.18, 5.14, 3.12, 0.54, size=10.5, color=MUTED,
                 line_spacing=1.10)

    add_text(slide, "3.87 px RMSE   ·   0.996 IoU   ·   2.45 px/frame",
             0.69, 6.35, 11.92, 0.42, size=16.4, color=BLUE, bold=True,
             align=PP_ALIGN.CENTER)
    add_notes(slide,
              "【约 35 秒】当前方法仍依赖可见、可辨认的屏幕边界。极低对比度、强反光和遮挡会削弱真实梯度，"
              "并产生更强的伪边缘。下一步应融合多种边界线索，为四条边分别估计置信度，并在连续失败时主动重初始化。"
              "但本次实验支持的结论已经清晰：让物理屏幕边框主导单应矩阵，可以同时获得更高的几何精度和更好的时间稳定性。"
              "整场汇报到这里约六分钟。")

    return prs


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT)
    print(f"Saved {len(prs.slides)} slides to {OUT}")


if __name__ == "__main__":
    main()
