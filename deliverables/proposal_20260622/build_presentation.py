"""Build the ECE4512 proposal presentation (.pptx).

English, 7 slides, 16:9, clean academic style matching the course example.
Re-runnable: regenerates ECE4512_Proposal_Presentation.pptx from assets/ + evidence/.

Usage:
    python deliverables/proposal_20260622/build_presentation.py
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- paths
HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = Path(os.environ.get("PPTX_OUT") or (HERE / "ECE4512_Proposal_Presentation.pptx"))

# ---------------------------------------------------------------- palette / type
INK = RGBColor(0x1A, 0x1A, 0x1A)        # near-black title/body
SUB = RGBColor(0x55, 0x5B, 0x66)        # muted subtitle/caption
ACCENT = RGBColor(0x1F, 0x57, 0xC4)     # blue accent (proposed method)
ACCENT_SOFT = RGBColor(0xE8, 0xEF, 0xFB)
WARM = RGBColor(0xC4, 0x3D, 0x2E)       # red (baseline / hard case)
CARD_BG = RGBColor(0xF4, 0xF6, 0xFA)
CARD_LINE = RGBColor(0xD3, 0xDA, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ARROW_GREY = RGBColor(0x9A, 0xA3, 0xB2)

FONT = "Calibri"
FONT_SCALE = 1.12  # global type scale; bump to enlarge all text uniformly

EMU_PER_IN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------- helpers
def _set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_autofit(tf):
    # prevent pptx from shrinking/auto-sizing; we control sizes explicitly
    tf.word_wrap = True


def add_text(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.05, space_after=4):
    """runs: list of paragraphs; each paragraph is a list of (text, size, bold, color, italic)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (text, size, bold, color, italic) in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size * FONT_SCALE)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    return box


def P(text, size=18, bold=False, color=INK, italic=False):
    """Single-run paragraph shortcut."""
    return [(text, size, bold, color, italic)]


def add_title(slide, text, *, kicker=None):
    if kicker:
        add_text(slide, Inches(0.6), Inches(0.34), Inches(12.1), Inches(0.3),
                 [P(kicker.upper(), 12.5, True, ACCENT)], space_after=0)
        ttop = Inches(0.62)
    else:
        ttop = Inches(0.42)
    add_text(slide, Inches(0.6), ttop, Inches(12.1), Inches(0.8),
             [P(text, 30, True, INK)], space_after=0)
    # accent underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), ttop + Inches(0.62),
                                  Inches(0.9), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


def add_picture_fit(slide, path, left, top, max_w, max_h, *, border=True):
    """Add picture scaled to fit max_w x max_h, centered in that box. Returns shape."""
    from PIL import Image
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 16, 9
    box_ratio = max_w / max_h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        w = max_w
        h = int(max_w / img_ratio)
    else:
        h = max_h
        w = int(max_h * img_ratio)
    l = left + (max_w - w) // 2
    t = top + (max_h - h) // 2
    pic = slide.shapes.add_picture(str(path), l, t, w, h)
    if border:
        pic.line.color.rgb = CARD_LINE
        pic.line.width = Pt(0.75)
    return pic


def add_caption(slide, left, top, width, text, *, align=PP_ALIGN.CENTER):
    add_text(slide, left, top, width, Inches(0.3), [P(text, 11.5, False, SUB, True)],
             align=align, space_after=0)


def add_card(slide, left, top, width, height, *, fill=CARD_BG, line=CARD_LINE,
             line_w=1.0, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    # tighter corner radius
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp


def add_chevron(slide, left, top, width, height, color=ARROW_GREY):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def footer(slide, page):
    add_text(slide, Inches(0.6), Inches(7.06), Inches(8.0), Inches(0.3),
             [P("ECE4512 Final Project Proposal · 2026", 9.5, False, SUB)], space_after=0)
    add_text(slide, Inches(11.4), Inches(7.06), Inches(1.3), Inches(0.3),
             [P(str(page), 9.5, False, SUB)], align=PP_ALIGN.RIGHT, space_after=0)


# ---------------------------------------------------------------- pipeline icons
# primitive shape helpers (coordinates in EMU) used to hand-draw the stage icons,
# styled after assets/application_pipeline.svg but in the deck's blue palette.
ICON_DARK = RGBColor(0x39, 0x3F, 0x4A)   # monitor bezel
WALL = RGBColor(0xEE, 0xF0, 0xF4)        # wall background
DESK = RGBColor(0xDC, 0xDF, 0xE6)        # desk background
JIT = RGBColor(0xE8, 0xA9, 0x9F)         # light red: raw jitter / rejected


def _rect(s, x, y, w, h, fill=None, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def _oval(s, cx, cy, r, fill, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(cx - r)), Emu(int(cy - r)), Emu(int(2 * r)), Emu(int(2 * r)))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def _seg(s, x1, y1, x2, y2, color, w=1.2, dash=None):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    if dash:
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
    return cn


def _polyline(s, pts, color, w, dash=None):
    for i in range(len(pts) - 1):
        _seg(s, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], color, w, dash)


def icon_input(s, ix, iy, iw, ih):
    """Full scene: wall/desk background with a small monitor — screen is one region."""
    _rect(s, ix, iy, iw, ih * 0.6, fill=WALL, line=CARD_LINE, lw=1.0)
    _rect(s, ix, iy + ih * 0.6, iw, ih * 0.4, fill=DESK, line=CARD_LINE, lw=1.0)
    mw, mh = iw * 0.4, ih * 0.4
    mx, my = ix + iw * 0.3, iy + ih * 0.16
    _rect(s, mx, my, mw, mh, fill=ICON_DARK)
    _rect(s, mx + mw * 0.12, my + mh * 0.14, mw * 0.76, mh * 0.62, fill=ACCENT_SOFT, line=ACCENT, lw=1.3)
    _rect(s, mx + mw * 0.42, my + mh, mw * 0.16, ih * 0.1, fill=ICON_DARK)
    _seg(s, mx + mw * 0.15, my + mh + ih * 0.1, mx + mw * 0.85, my + mh + ih * 0.1, ICON_DARK, 2.0)


def icon_init(s, ix, iy, iw, ih):
    """Tilted quadrilateral with four detected corners."""
    _rect(s, ix, iy, iw, ih, fill=RGBColor(0xF8, 0xF9, 0xFB), line=CARD_LINE, lw=1.0)
    cs = [(0.20, 0.30), (0.80, 0.20), (0.86, 0.74), (0.16, 0.66)]
    pts = [(ix + iw * fx, iy + ih * fy) for fx, fy in cs]
    for i in range(4):
        x1, y1 = pts[i]; x2, y2 = pts[(i + 1) % 4]
        _seg(s, x1, y1, x2, y2, ACCENT, 2.0)
    r = int(min(iw, ih) * 0.06)
    for x, y in pts:
        _oval(s, x, y, r, ACCENT)


def icon_track(s, ix, iy, iw, ih):
    """Three mini frames t0,t1,t2 with a feature point tracked along a path."""
    fw, fh, ys = iw * 0.26, ih * 0.58, iy + ih * 0.13
    centers = []
    for fx in (0.04, 0.37, 0.70):
        fx0 = ix + iw * fx
        _rect(s, fx0, ys, fw, fh, fill=WHITE, line=CARD_LINE, lw=1.0)
        _rect(s, fx0 + fw * 0.16, ys + fh * 0.16, fw * 0.68, fh * 0.44, fill=ACCENT_SOFT, line=ACCENT, lw=1.1)
        centers.append((fx0 + fw * 0.5, ys + fh * 0.4))
    _polyline(s, centers, WARM, 1.6)
    for cx, cy in centers:
        _oval(s, cx, cy, int(min(iw, ih) * 0.045), ACCENT)


def icon_ransac(s, ix, iy, iw, ih):
    """Screen plane with blue inliers (kept) and red outliers (content motion, rejected)."""
    _rect(s, ix, iy, iw, ih, fill=WHITE, line=CARD_LINE, lw=1.0)
    _rect(s, ix + iw * 0.18, iy + ih * 0.24, iw * 0.64, ih * 0.4, fill=ACCENT_SOFT, line=ACCENT, lw=1.8)
    r = int(min(iw, ih) * 0.05)
    ins = [(0.24, 0.28), (0.76, 0.26), (0.5, 0.62)]
    ipts = [(ix + iw * fx, iy + ih * fy) for fx, fy in ins]
    _seg(s, ipts[0][0], ipts[0][1], ipts[1][0], ipts[1][1], ACCENT, 1.4)
    for x, y in ipts:
        _oval(s, x, y, r, ACCENT)
    for fx, fy in ((0.42, 0.46), (0.6, 0.5)):
        x, y = ix + iw * fx, iy + ih * fy
        _seg(s, x, y, x + iw * 0.13, y + ih * 0.26, WARM, 1.3, dash="dash")
        _oval(s, x, y, r, WARM)


def icon_filter(s, ix, iy, iw, ih):
    """Axes with a jittery raw trajectory (red) and the smoothed one (blue)."""
    ax0, ay0 = ix + iw * 0.12, iy + ih * 0.84
    _seg(s, ax0, iy + ih * 0.1, ax0, ay0, CARD_LINE, 1.2)
    _seg(s, ax0, ay0, ix + iw * 0.92, ay0, CARD_LINE, 1.2)
    jit = [(0.14, 0.5), (0.26, 0.22), (0.38, 0.74), (0.5, 0.36), (0.62, 0.68), (0.76, 0.42), (0.9, 0.54)]
    _polyline(s, [(ix + iw * fx, iy + ih * fy) for fx, fy in jit], JIT, 1.6)
    sm = [(0.14, 0.56), (0.32, 0.5), (0.5, 0.47), (0.7, 0.45), (0.9, 0.44)]
    _polyline(s, [(ix + iw * fx, iy + ih * fy) for fx, fy in sm], ACCENT, 2.6)


def icon_output(s, ix, iy, iw, ih):
    """Rectified front-facing screen with a 16:9 ruler."""
    sx, sy, sw, sh = ix + iw * 0.16, iy + ih * 0.1, iw * 0.68, ih * 0.5
    _rect(s, sx, sy, sw, sh, fill=WHITE, line=ACCENT, lw=2.2)
    _rect(s, sx + sw * 0.1, sy + sh * 0.22, sw * 0.62, sh * 0.12, fill=CARD_LINE)
    _rect(s, sx + sw * 0.1, sy + sh * 0.48, sw * 0.42, sh * 0.12, fill=CARD_LINE)
    ry = iy + ih * 0.82
    _seg(s, ix + iw * 0.14, ry, ix + iw * 0.86, ry, ACCENT, 2.0)
    for fx in (0.18, 0.5, 0.82):
        _seg(s, ix + iw * fx, ry - ih * 0.05, ix + iw * fx, ry + ih * 0.05, ACCENT, 1.6)


ICONS = [icon_input, icon_init, icon_track, icon_ransac, icon_filter, icon_output]


# ---------------------------------------------------------------- deck
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    _set_bg(s)
    return s


# ===== Slide 1 : Title + Motivation =========================================
s = new_slide()
# left text column
add_text(s, Inches(0.7), Inches(0.55), Inches(7.6), Inches(0.3),
         [P("ECE4512 FINAL PROJECT PROPOSAL · 2026", 13, True, ACCENT)], space_after=0)
add_text(s, Inches(0.7), Inches(1.0), Inches(7.7), Inches(2.0),
         [P("Screen Capture Rectification and", 33, True, INK),
          P("Temporal Stabilization for", 33, True, INK),
          P("Real-world Captured-screen Videos", 33, True, INK)],
         line_spacing=1.02, space_after=2)

add_text(s, Inches(0.72), Inches(3.05), Inches(7.6), Inches(0.7),
         [[("Rongshuo Wen", 14.5, True, INK, False), ("  124020369      ", 13, False, SUB, False),
           ("Bihua Wen", 14.5, True, INK, False), ("  124090670", 13, False, SUB, False)],
          [("Mingrui Liu", 14.5, True, INK, False), ("  124090375", 13, False, SUB, False)]],
         line_spacing=1.15, space_after=2)

# motivation card
add_card(s, Inches(0.7), Inches(4.05), Inches(7.7), Inches(2.55), fill=ACCENT_SOFT,
         line=None)
add_text(s, Inches(0.95), Inches(4.25), Inches(7.2), Inches(2.2),
         [[("Motivation.  ", 14.5, True, ACCENT, False),
           ("Screen demoiréing and screen-image restoration show captured-screen "
            "content is a real image-processing problem — yet most datasets start "
            "after the screen is already cropped, aligned, or paired with a clean "
            "reference.", 14.5, False, INK, False)],
          [("Real phone-captured video is far less controlled: background clutter, "
            "perspective tilt, hand-held shake, weak borders, glare, moiré, and "
            "dynamic on-screen content. We tackle this ", 14.5, False, INK, False),
           ("missing geometric preprocessing step", 14.5, True, INK, False),
           (".", 14.5, False, INK, False)]],
         line_spacing=1.12, space_after=8)

# right column: our stage (top) + the downstream task it enables (bottom)
add_picture_fit(s, ASSETS / "comparison_4s.jpg", Inches(8.65), Inches(1.7),
                Inches(4.25), Inches(1.5))
add_caption(s, Inches(8.65), Inches(3.12), Inches(4.25),
            "Our stage:  handheld capture  →  rectified screen")
add_picture_fit(s, ASSETS / "downstream_demoire.png", Inches(8.65), Inches(3.7),
                Inches(4.25), Inches(1.55))
add_caption(s, Inches(8.65), Inches(5.28), Inches(4.25),
            "Downstream enabled:  video demoiréing on the stable output")
add_text(s, Inches(8.65), Inches(5.66), Inches(4.25), Inches(1.0),
         [P("This project is the geometry stage BEFORE restoration: locate, "
            "rectify, and stabilize the screen.", 12.5, False, SUB, True)],
         align=PP_ALIGN.CENTER, line_spacing=1.1, space_after=0)
footer(s, 1)

# ===== Slide 2 : Problem & Goal =============================================
s = new_slide()
add_title(s, "Problem & Goal", kicker="The gap we fill")

# left: the gap
add_text(s, Inches(0.65), Inches(1.55), Inches(6.0), Inches(0.4),
         [P("The gap", 16, True, ACCENT)], space_after=2)
add_text(s, Inches(0.65), Inches(1.98), Inches(6.0), Inches(2.4),
         [P("▸  Downstream restoration assumes the screen is already cropped & "
            "aligned.", 14, False, INK),
          P("▸  Real handheld video has off-screen background, perspective tilt, "
            "and frame-to-frame jitter.", 14, False, INK),
          P("▸  Weak borders, glare and moiré make a single frame-by-frame "
            "detector unstable.", 14, False, INK)],
         line_spacing=1.12, space_after=9)

# goal card
add_card(s, Inches(0.65), Inches(4.35), Inches(6.0), Inches(2.35), fill=CARD_BG)
add_text(s, Inches(0.9), Inches(4.55), Inches(5.55), Inches(2.0),
         [[("Goal.  ", 15, True, ACCENT, False),
           ("Continuously estimate the screen plane and render its content as a "
            "front-facing video that:", 14, False, INK, False)],
          P("•  suppresses non-screen background & perspective", 13.5, False, INK),
          P("•  reduces hand-held jitter while keeping true content motion", 13.5, False, INK),
          P("•  preserves physical aspect ratio  (default 1920×1080)", 13.5, False, INK)],
         line_spacing=1.12, space_after=5)

# right: visual (terminal capture example: raw → normalized)
add_picture_fit(s, ASSETS / "comparison_terminal.jpg", Inches(7.0), Inches(1.7),
                Inches(5.85), Inches(3.95))
add_caption(s, Inches(7.0), Inches(5.7), Inches(5.85),
            "Same frame: raw handheld input  →  rectified, stable screen output")
footer(s, 2)

# ===== Slide 3 : Method - Pipeline ==========================================
s = new_slide()
add_title(s, "Method — Geometric Pipeline", kicker="Approach overview")

# pipeline stages, each with a hand-drawn icon (style after application_pipeline.svg)
stages = [
    ("INPUT", "Handheld video", "screen is one region\nin a full scene"),
    ("STAGE 1", "Init plane", "detect / set four\ncorners at t₀"),
    ("STAGE 2", "Track features", "LK tracks p₀ → pₜ\non the t₀ plane"),
    ("STAGE 3", "Estimate Hₜ", "RANSAC screen-plane\nhomography"),
    ("STAGE 4", "Filter trajectory", "gate, interpolate\n& smooth Hₜ"),
    ("OUTPUT", "Rectified screen", "stable 16:9 video\nfor later modules"),
]
n = len(stages)
top = Inches(1.62)
card_h = Inches(2.5)
left0 = Inches(0.5)
total_w = Inches(12.33)
gap = Inches(0.2)
card_w = Emu(int((total_w - gap * (n - 1)) / n))
icon_top = int(top + Inches(0.48))
icon_h = int(Inches(1.02))
for i, (kick, title, body) in enumerate(stages):
    l = Emu(int(left0 + i * (card_w + gap)))
    is_io = kick in ("INPUT", "OUTPUT")
    add_card(s, l, top, card_w, card_h,
             fill=ACCENT_SOFT if is_io else WHITE,
             line=ACCENT if is_io else CARD_LINE,
             line_w=1.6 if is_io else 1.0)
    add_text(s, l + Inches(0.08), top + Inches(0.13), card_w - Inches(0.16), Inches(0.3),
             [P(kick, 10, True, ACCENT if is_io else SUB)], align=PP_ALIGN.CENTER, space_after=0)
    ICONS[i](s, int(l + Inches(0.2)), icon_top, int(card_w - Inches(0.4)), icon_h)
    add_text(s, l + Inches(0.05), top + Inches(1.62), card_w - Inches(0.1), Inches(0.35),
             [P(title, 12.5, True, INK)], align=PP_ALIGN.CENTER, space_after=0)
    add_text(s, l + Inches(0.05), top + Inches(2.0), card_w - Inches(0.1), Inches(0.5),
             [P(body, 9.5, False, SUB)], align=PP_ALIGN.CENTER, line_spacing=1.0, space_after=0)
    if i < n - 1:
        ax = Emu(int(l + card_w + Inches(0.012)))
        add_chevron(s, ax, top + Inches(0.86), Inches(0.18), Inches(0.32))

# downstream bar
add_card(s, Inches(0.6), Inches(4.34), Inches(12.13), Inches(0.6), fill=ACCENT_SOFT, line=None)
add_text(s, Inches(0.6), Inches(4.48), Inches(12.13), Inches(0.32),
         [[("Downstream after this stage:    ", 12.5, True, ACCENT, False),
           ("video demoiréing    ·    OCR    ·    archival    ·    manual review",
            12.5, False, INK, True)]],
         align=PP_ALIGN.CENTER, space_after=0)

# real example under the schematic pipeline
add_picture_fit(s, ASSETS / "method_tracking_diagram.png", Inches(4.2), Inches(5.06),
                Inches(4.9), Inches(1.68))
add_caption(s, Inches(3.7), Inches(6.78), Inches(5.9),
            "Real example: tilted frames → corners + homography → upright, plane-locked frames")
footer(s, 3)

# ===== Slide 4 : Screen vs Content Motion ===================================
s = new_slide()
add_title(s, "Method — Screen Motion vs. Content Motion", kicker="Why it stays stable")

add_text(s, Inches(0.65), Inches(1.5), Inches(12.0), Inches(0.5),
         [[("Core distinction:  ", 15, True, ACCENT, False),
           ("estimate the homography from the physical screen border, and treat motion "
            "inside the screen as content — not camera — motion.", 15, False, INK, False)]],
         line_spacing=1.1, space_after=0)

cols = [
    ("1 · Border drives geometry",
     ["Edge filtering + LSD / Hough locate the four screen borders.",
      "The main homography is anchored to this physical boundary, not to texture inside the screen."]),
    ("2 · Inner points = check only",
     ["Inner Lucas–Kanade points serve as a consistency signal.",
      "If inner motion conflicts with border motion under RANSAC, it is labelled screen-content "
      "motion and excluded from H estimation."]),
    ("3 · Robust fallback",
     ["Low border confidence, low inlier ratio, or an invalid quadrilateral triggers re-detection.",
      "If recovery fails, freeze the last valid Hₜ and flag invalid regions; corner trajectory "
      "is finally smoothed in time."]),
]
cw = Inches(3.95)
cgap = Inches(0.27)
ctop = Inches(2.35)
ch = Inches(2.8)
for i, (head, bullets) in enumerate(cols):
    l = Emu(int(Inches(0.65) + i * (cw + cgap)))
    add_card(s, l, ctop, cw, ch, fill=CARD_BG)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, ctop, cw, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(s, l + Inches(0.25), ctop + Inches(0.28), cw - Inches(0.5), Inches(0.6),
             [P(head, 14.5, True, INK)], line_spacing=1.05, space_after=0)
    paras = [[("•  ", 13, False, ACCENT, False), (b, 13, False, INK, False)] for b in bullets]
    add_text(s, l + Inches(0.25), ctop + Inches(1.05), cw - Inches(0.5), ch - Inches(1.2),
             paras, line_spacing=1.13, space_after=10)

# bottom: normalized time strip (left) + result note (right)
add_caption(s, Inches(0.65), Inches(5.18), Inches(6.2),
            "Normalized output over time  (content moves, frame stays aligned)",
            align=PP_ALIGN.LEFT)
add_picture_fit(s, ASSETS / "normalized_time_strip.png", Inches(0.65), Inches(5.45),
                Inches(6.2), Inches(1.5))
add_text(s, Inches(7.15), Inches(5.55), Inches(5.6), Inches(1.45),
         [[("Result:  ", 13.5, True, ACCENT, False),
           ("in-screen content changes over time, yet each frame stays rectified "
            "and plane-aligned — screen motion removed, content motion preserved.",
            13.5, False, INK, True)]],
         line_spacing=1.18, space_after=0)
footer(s, 4)

# ===== Slide 5 : Dataset & Experiment =======================================
s = new_slide()
add_title(s, "Dataset & Experiment", kicker="Self-collected evaluation set")

add_text(s, Inches(0.65), Inches(1.5), Inches(12.0), Inches(0.4),
         [[("5 scenario classes × 10 clips × ~5 s  =  ", 15, True, INK, False),
           ("50 real captured-screen clips", 15, True, ACCENT, False),
           (".  Classes are designed around failure modes.", 15, False, INK, False)]],
         space_after=0)

classes = [
    ("CLASS 1", "Static page", "documents, browser UI,\nmostly fixed content"),
    ("CLASS 2", "Scrolling page", "page motion, cursor,\nlocal UI changes"),
    ("CLASS 3", "In-screen video", "large dynamic regions\ninside the display"),
    ("CLASS 4", "PPT / weak border", "low texture, pale UI,\nunclear screen edges"),
    ("CLASS 5", "4K / moiré / glare", "hard cases linking to\nthe restoration pipeline"),
]
n = len(classes)
top = Inches(2.25)
ch = Inches(2.35)
left0 = Inches(0.65)
total_w = Inches(12.05)
gap = Inches(0.22)
cw = Emu(int((total_w - gap * (n - 1)) / n))
for i, (kick, title, body) in enumerate(classes):
    l = Emu(int(left0 + i * (cw + gap)))
    hard = i == n - 1
    add_card(s, l, top, cw, ch, fill=RGBColor(0xFB, 0xEC, 0xEA) if hard else CARD_BG,
             line=WARM if hard else CARD_LINE, line_w=1.4 if hard else 1.0)
    add_text(s, l + Inches(0.05), top + Inches(0.18), cw - Inches(0.1), Inches(0.3),
             [P(kick, 10.5, True, WARM if hard else ACCENT)], align=PP_ALIGN.CENTER, space_after=0)
    add_text(s, l + Inches(0.05), top + Inches(0.6), cw - Inches(0.1), Inches(0.6),
             [P(title, 13.5, True, INK)], align=PP_ALIGN.CENTER, line_spacing=1.0, space_after=0)
    add_text(s, l + Inches(0.05), top + Inches(1.35), cw - Inches(0.1), Inches(0.7),
             [P(body, 10, False, SUB)], align=PP_ALIGN.CENTER, line_spacing=1.0, space_after=0)
    add_text(s, l + Inches(0.05), top + Inches(2.02), cw - Inches(0.1), Inches(0.28),
             [P("10 clips × 5 s", 10, True, INK)], align=PP_ALIGN.CENTER, space_after=0)

add_card(s, Inches(0.65), Inches(5.05), Inches(12.05), Inches(1.55), fill=ACCENT_SOFT, line=None)
add_text(s, Inches(0.95), Inches(5.25), Inches(11.5), Inches(1.2),
         [[("Annotation.  ", 14.5, True, ACCENT, False),
           ("Selected key frames from each class are manually labelled with the four "
            "screen corners — connecting quantitative metrics to visual comparison, "
            "without requiring any downstream demoiréing model.", 14.5, False, INK, False)]],
         line_spacing=1.18, space_after=0)
footer(s, 5)

# ===== Slide 6 : Evaluation Metrics =========================================
s = new_slide()
add_title(s, "Evaluation Metrics", kicker="What we measure")

groups = [
    ("Geometric accuracy", ACCENT,
     ["Corner error on annotated key frames",
      "Quadrilateral IoU",
      "Aspect-ratio error"],
     "Is the screen located & rectified correctly?"),
    ("Temporal stability", ACCENT,
     ["Residual adjacent-frame translation (px)",
      "Residual rotation (deg) & scale change",
      "Reported as p95 of the normalized video"],
     "How much jitter remains after normalization?"),
    ("Signal preservation", ACCENT,
     ["Mean gradient magnitude on texture patches",
      "Edge-preservation index (same screen scale)",
      "2D FFT grid orthogonality for moiré cases"],
     "Does rectification keep high-frequency detail?"),
]
cw = Inches(3.95)
cgap = Inches(0.27)
ctop = Inches(1.65)
ch = Inches(4.05)
for i, (head, col, bullets, foot) in enumerate(groups):
    l = Emu(int(Inches(0.65) + i * (cw + cgap)))
    add_card(s, l, ctop, cw, ch, fill=CARD_BG)
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.25), ctop + Inches(0.28),
                               Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background()
    badge.shadow.inherit = False
    bt = badge.text_frame; bt.word_wrap = False
    bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = str(i + 1)
    br.font.name = FONT; br.font.size = Pt(18 * FONT_SCALE); br.font.bold = True; br.font.color.rgb = WHITE
    add_text(s, l + Inches(0.9), ctop + Inches(0.32), cw - Inches(1.05), Inches(0.7),
             [P(head, 15, True, INK)], line_spacing=0.95, space_after=0)
    paras = [[("•  ", 13, False, col, False), (b, 13, False, INK, False)] for b in bullets]
    add_text(s, l + Inches(0.27), ctop + Inches(1.3), cw - Inches(0.5), Inches(1.9),
             paras, line_spacing=1.18, space_after=10)
    add_text(s, l + Inches(0.27), ctop + Inches(3.35), cw - Inches(0.5), Inches(0.6),
             [P(foot, 11.5, False, SUB, True)], line_spacing=1.05, space_after=0)

add_text(s, Inches(0.65), Inches(5.95), Inches(12.0), Inches(0.5),
         [P("These measures tie geometric correctness, video stability and signal "
            "preservation together — with no dependence on a downstream restoration model.",
            13, False, SUB, True)],
         align=PP_ALIGN.CENTER, line_spacing=1.1, space_after=0)
footer(s, 6)

# ===== Slide 7 : Initial Results + Timeline =================================
s = new_slide()
add_title(s, "Initial Results & Timeline", kicker="Border-guided tracking wins")

# left: ablation chart
add_picture_fit(s, ASSETS / "method_ablation_translation_p95.png", Inches(0.55), Inches(1.55),
                Inches(7.0), Inches(3.55), border=False)
add_caption(s, Inches(0.55), Inches(5.12), Inches(7.0),
            "Same input video · residual translation p95 over the last 2 s (lower = steadier)")

# right: result highlight + comparison of three methods
add_card(s, Inches(7.85), Inches(1.55), Inches(4.95), Inches(2.05), fill=ACCENT_SOFT, line=None)
add_text(s, Inches(8.1), Inches(1.72), Inches(4.5), Inches(1.8),
         [[("0.118 px", 30, True, ACCENT, False),
           ("  residual jitter", 14, False, INK, False)],
          [("Reference-plane (border-guided) tracking cuts residual motion ", 13, False, INK, False),
           ("~16×", 13, True, ACCENT, False),
           (" vs. per-frame detection (1.927) and optical-flow tracking (1.929).",
            13, False, INK, False)]],
         line_spacing=1.12, space_after=6)

add_text(s, Inches(7.85), Inches(3.75), Inches(5.0), Inches(0.35),
         [P("Final report will compare three strategies:", 13, True, INK)], space_after=2)
comp = [
    ("Per-frame detection", "re-detects each frame — jittery", WARM),
    ("Optical-flow tracking", "drifts with on-screen content", WARM),
    ("Border-guided (ours)", "stable, content-aware", ACCENT),
]
yt = Inches(4.15)
for j, (name, desc, c) in enumerate(comp):
    y = Emu(int(yt + j * Inches(0.5)))
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.9), y + Inches(0.06), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = c; dot.line.fill.background(); dot.shadow.inherit = False
    add_text(s, Inches(8.2), y, Inches(4.6), Inches(0.45),
             [[(name + " — ", 12.5, True, INK, False), (desc, 12.5, False, SUB, True)]],
             space_after=0)

# timeline bar at bottom
add_text(s, Inches(0.6), Inches(5.75), Inches(6.0), Inches(0.3),
         [P("Timeline", 14, True, ACCENT)], space_after=0)
milestones = [
    ("Jun 22–24", "Proposal & presentation"),
    ("Jun 25–30", "50-clip set + key-frame labels"),
    ("Jul 1–7", "Method ablation & metrics"),
    ("Jul 8–12", "Visual comparison & report"),
    ("Jul 13–15", "Code, samples, final talk"),
]
n = len(milestones)
top = Inches(6.15)
left0 = Inches(0.6)
total_w = Inches(12.15)
gap = Inches(0.16)
mw = Emu(int((total_w - gap * (n - 1)) / n))
for i, (when, what) in enumerate(milestones):
    l = Emu(int(left0 + i * (mw + gap)))
    done = i == 0
    add_card(s, l, top, mw, Inches(0.82), fill=ACCENT_SOFT if done else CARD_BG,
             line=ACCENT if done else CARD_LINE, line_w=1.3 if done else 1.0)
    add_text(s, l + Inches(0.1), top + Inches(0.08), mw - Inches(0.2), Inches(0.28),
             [P(when, 11, True, ACCENT if done else INK)], align=PP_ALIGN.CENTER, space_after=0)
    add_text(s, l + Inches(0.08), top + Inches(0.36), mw - Inches(0.16), Inches(0.42),
             [P(what, 9.5, False, SUB)], align=PP_ALIGN.CENTER, line_spacing=0.95, space_after=0)
footer(s, 7)

# ---------------------------------------------------------------- save
prs.save(str(OUT))
print(f"Saved: {OUT}  ({len(prs.slides._sldIdLst)} slides)")
